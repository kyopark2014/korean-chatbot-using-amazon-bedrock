import json
import boto3
import os
import time
import datetime
from io import BytesIO
import PyPDF2
import csv
import traceback
import re
import base64
import datetime
import requests
import docx

from urllib import parse
from botocore.config import Config
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferWindowMemory

from langchain_community.docstore.document import Document
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_aws import BedrockEmbeddings
from langchain_community.retrievers import AmazonKendraRetriever
from multiprocessing import Process, Pipe
from googleapiclient.discovery import build
from opensearchpy import OpenSearch
from langchain_core.prompts import PromptTemplate

from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock

from langchain.agents import tool
from langchain.agents import AgentExecutor, create_react_agent
from bs4 import BeautifulSoup
from pytz import timezone
from langchain_community.tools.tavily_search import TavilySearchResults

from typing import TypedDict, Annotated, Sequence, List, Union
from langchain_core.agents import AgentAction, AgentFinish
from langchain_core.messages import BaseMessage, HumanMessage, AIMessage, SystemMessage
from langgraph.prebuilt.tool_executor import ToolExecutor
from langgraph.graph import START, END, StateGraph
from langgraph.prebuilt import ToolNode
from langgraph.graph.message import add_messages
from langgraph.prebuilt import tools_condition
from langchain_core.pydantic_v1 import BaseModel, Field
from typing import Literal
from langchain_aws import AmazonKnowledgeBasesRetriever

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
callLogTableName = os.environ.get('callLogTableName')
kendra_region = os.environ.get('kendra_region', 'us-west-2')
LLM_for_chat = json.loads(os.environ.get('LLM_for_chat'))
LLM_for_multimodal= json.loads(os.environ.get('LLM_for_multimodal'))
LLM_embedding = json.loads(os.environ.get('LLM_embedding'))
priority_search_embedding = json.loads(os.environ.get('priority_search_embedding'))
selected_chat = 0
selected_multimodal = 0
selected_embedding = 0
selected_ps_embedding = 0
rag_method = os.environ.get('rag_method', 'RetrievalPrompt') # RetrievalPrompt, RetrievalQA, ConversationalRetrievalChain
separated_chat_history = os.environ.get('separated_chat_history')
enalbeParentDocumentRetrival = os.environ.get('enalbeParentDocumentRetrival')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'
speech_prefix = 'speech/'

useParallelRAG = os.environ.get('useParallelRAG', 'true')
kendraIndex = os.environ.get('kendraIndex')
kendra_method = os.environ.get('kendraMethod')
roleArn = os.environ.get('roleArn')
top_k = int(os.environ.get('numberOfRelevantDocs'))
capabilities = json.loads(os.environ.get('capabilities'))
print('capabilities: ', capabilities)
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = os.environ.get('speech_generation')
history_length = 0
token_counter_history = 0
allowDualSearch = os.environ.get('allowDualSearch')
allowDualSearchWithMulipleProcessing = True
enableHybridSearch = os.environ.get('enableHybridSearch')
useParrelWebSearch = True
useEnhancedSearch = True

minDocSimilarity = 200
minCodeSimilarity = 300
projectName = os.environ.get('projectName')

flow_id = os.environ.get('flow_id')
flow_alias = os.environ.get('flow_alias')

rag_flow_id = os.environ.get('rag_flow_id')
rag_flow_alias = os.environ.get('rag_flow_alias')

reference_docs = []

# google search api
googleApiSecret = os.environ.get('googleApiSecret')
secretsmanager = boto3.client('secretsmanager')
try:
    get_secret_value_response = secretsmanager.get_secret_value(
        SecretId=googleApiSecret
    )
    #print('get_secret_value_response: ', get_secret_value_response)
    secret = json.loads(get_secret_value_response['SecretString'])
    #print('secret: ', secret)
    google_api_key = secret['google_api_key']
    google_cse_id = secret['google_cse_id']
    #print('google_cse_id: ', google_cse_id)    

except Exception as e:
    raise e

# api key to get weather information in agent
try:
    get_weather_api_secret = secretsmanager.get_secret_value(
        SecretId=f"openweathermap-{projectName}"
    )
    #print('get_weather_api_secret: ', get_weather_api_secret)
    secret = json.loads(get_weather_api_secret['SecretString'])
    #print('secret: ', secret)
    weather_api_key = secret['weather_api_key']

except Exception as e:
    raise e
   
# api key to use LangSmith
langsmith_api_key = ""
try:
    get_langsmith_api_secret = secretsmanager.get_secret_value(
        SecretId=f"langsmithapikey-{projectName}"
    )
    #print('get_langsmith_api_secret: ', get_langsmith_api_secret)
    secret = json.loads(get_langsmith_api_secret['SecretString'])
    #print('secret: ', secret)
    langsmith_api_key = secret['langsmith_api_key']
    langchain_project = secret['langchain_project']
except Exception as e:
    raise e

if langsmith_api_key:
    os.environ["LANGCHAIN_API_KEY"] = langsmith_api_key
    os.environ["LANGCHAIN_TRACING_V2"] = "true"
    os.environ["LANGCHAIN_PROJECT"] = langchain_project
    
# api key to use Tavily Search
tavily_api_key = ""
try:
    get_tavily_api_secret = secretsmanager.get_secret_value(
        SecretId=f"tavilyapikey-{projectName}"
    )
    #print('get_tavily_api_secret: ', get_tavily_api_secret)
    secret = json.loads(get_tavily_api_secret['SecretString'])
    #print('secret: ', secret)
    tavily_api_key = secret['tavily_api_key']
except Exception as e: 
    raise e

if tavily_api_key:
    os.environ["TAVILY_API_KEY"] = tavily_api_key
    
# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"

map_chain = dict() 

# Multi-LLM
def get_chat():
    global selected_chat
    profile = LLM_for_chat[selected_chat]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'LLM: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_chat = selected_chat + 1
    if selected_chat == len(LLM_for_chat):
        selected_chat = 0
    
    return chat

def get_multi_region_chat(models, selected):
    profile = models[selected]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'selected_chat: {selected}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    return chat

def get_multimodal():
    global selected_multimodal
    
    profile = LLM_for_multimodal[selected_multimodal]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'LLM: {selected_multimodal}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    multimodal = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_multimodal = selected_multimodal + 1
    if selected_multimodal == len(LLM_for_multimodal):
        selected_multimodal = 0
    
    return multimodal
    
def get_embedding():
    global selected_embedding
    profile = LLM_embedding[selected_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_embedding: {selected_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_embedding = selected_embedding + 1
    if selected_embedding == len(LLM_embedding):
        selected_embedding = 0
    
    return bedrock_embedding

def get_ps_embedding():
    global selected_ps_embedding
    profile = priority_search_embedding[selected_ps_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_ps_embedding: {selected_ps_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_ps_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_ps_embedding = selected_ps_embedding + 1
    if selected_ps_embedding == len(priority_search_embedding):
        selected_ps_embedding = 0
    
    return bedrock_ps_embedding
    
def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)
    
os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def general_conversation(connectionId, requestId, chat, query):
    global time_for_inference, history_length, token_counter_history    
    time_for_inference = history_length = token_counter_history = 0
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
    
    if isKorean(query)==True :
        system = (
            "다음의 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다."
        )
    else: 
        system = (
            "Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor."
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "history": history,
                "input": query,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
                            
        usage = stream.response_metadata['usage']
        print('prompt_tokens: ', usage['prompt_tokens'])
        print('completion_tokens: ', usage['completion_tokens'])
        print('total_tokens: ', usage['total_tokens'])
        msg = stream.content
        # print('msg: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
        
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
        
    return msg

def traslation(chat, text, input_language, output_language):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        # print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def translate_text(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    if isKorean(text)==False :
        input_language = "English"
        output_language = "Korean"
    else:
        input_language = "Korean"
        output_language = "English"
                        
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def check_grammer(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Find the error in the sentence and explain it, and add the corrected sentence at the end of your answer."
        )
        
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        msg = result.content
        print('result of grammer correction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg

def extract_sentiment(chat, text):
    if isKorean(text)==True:
        system = (
            """아래의 <example> review와 Extracted Topic and sentiment 인 <result>가 있습니다.
            <example>
            객실은 작지만 깨끗하고 편안합니다. 프론트 데스크는 정말 분주했고 체크인 줄도 길었지만, 직원들은 프로페셔널하고 매우 유쾌하게 각 사람을 응대했습니다. 우리는 다시 거기에 머물것입니다.
            </example>
            <result>
            청소: 긍정적, 
            서비스: 긍정적
            </result>

            아래의 <review>에 대해서 위의 <result> 예시처럼 Extracted Topic and sentiment 을 만들어 주세요."""
        )
    else: 
        system = (
            """Here is <example> review and extracted topics and sentiments as <result>.

            <example>
            The room was small but clean and comfortable. The front desk was really busy and the check-in line was long, but the staff were professional and very pleasant with each person they helped. We will stay there again.
            </example>

            <result>
            Cleanliness: Positive, 
            Service: Positive
            </result>"""
        )
        
    human = "<review>{text}</review>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )        
        msg = result.content                
        print('result of sentiment extraction: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return msg

def extract_information(chat, text):
    if isKorean(text)==True:
        system = (
            """다음 텍스트에서 이메일 주소를 정확하게 복사하여 한 줄에 하나씩 적어주세요. 입력 텍스트에 정확하게 쓰여있는 이메일 주소만 적어주세요. 텍스트에 이메일 주소가 없다면, "N/A"라고 적어주세요. 또한 결과는 <result> tag를 붙여주세요."""
        )
    else: 
        system = (
            """Please precisely copy any email addresses from the following text and then write them, one per line.  Only write an email address if it's precisely spelled out in the input text. If there are no email addresses in the text, write "N/A".  Do not say anything else.  Put it in <result> tags."""
        )
        
    human = "<text>{text}</text>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )        
        output = result.content        
        msg = output[output.find('<result>')+8:len(output)-9] # remove <result> 
        
        print('result of information extraction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return msg

def remove_pii(chat, text):
    if isKorean(text)==True:
        system = (
            """아래의 <text>에서 개인식별정보(PII)를 모두 제거하여 외부 계약자와 안전하게 공유할 수 있도록 합니다. 이름, 전화번호, 주소, 이메일을 XXX로 대체합니다. 또한 결과는 <result> tag를 붙여주세요."""
        )
    else: 
        system = (
            """We want to de-identify some text by removing all personally identifiable information from this text so that it can be shared safely with external contractors.
            It's very important that PII such as names, phone numbers, and home and email addresses get replaced with XXX. Put it in <result> tags."""
        )
        
    human = "<text>{text}</text>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )        
        output = result.content        
        msg = output[output.find('<result>')+8:len(output)-9] # remove <result> 
        
        print('result of removing PII : ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return msg

def do_step_by_step(chat, text):
    if isKorean(text)==True:
        system = (
            """다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. 아래 문맥(context)을 참조했음에도 답을 알 수 없다면, 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다.

            Assistant: 단계별로 생각할까요?

            Human: 예, 그렇게하세요."""
        )
    else: 
        system = (
            """Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor. 
            
            Assistant: Can I think step by step?

            Human: Yes, please do."""
        )
        
    human = "<text>{text}</text>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )        
        msg = result.content        
        
        print('result of sentiment extraction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return msg

def extract_timestamp(chat, text):
    system = (
        """Human: 아래의 <text>는 시간을 포함한 텍스트입니다. 친절한 AI Assistant로서 시간을 추출하여 아래를 참조하여 <example>과 같이 정리해주세요.
            
        - 년도를 추출해서 <year>/<year>로 넣을것 
        - 월을 추출해서 <month>/<month>로 넣을것
        - 일을 추출해서 <day>/<day>로 넣을것
        - 시간을 추출해서 24H으로 정리해서 <hour>/<hour>에 넣을것
        - 분을 추출해서 <minute>/<minute>로 넣을것

        이때의 예제는 아래와 같습니다.
        <example>
        2022년 11월 3일 18시 26분
        </example>
        <result>
            <year>2022</year>
            <month>11</month>
            <day>03</day>
            <hour>18</hour>
            <minute>26</minute>
        </result>

        결과에 개행문자인 "\n"과 글자 수와 같은 부가정보는 절대 포함하지 마세요."""
    )    
        
    human = "<text>{text}</text>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )        
        output = result.content        
        msg = output[output.find('<result>')+8:len(output)-9] # remove <result> 
        
        print('result of sentiment extraction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return msg

def use_multimodal(img_base64, query):    
    multimodal = get_multimodal()
    
    if query == "":
        query = "그림에 대해 상세히 설명해줘."
    
    messages = [
        SystemMessage(content="답변은 500자 이내의 한국어로 설명해주세요."),
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = multimodal.invoke(messages)
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def get_prompt_template(query, conv_type, rag_type):    
    if isKorean(query):
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

            <history>
            {history}
            </history>            

            <question>            
            {input}
            </question>
            
            Assistant:"""

        elif conv_type=='qa':  
            # for RAG, context and question
            #prompt_template = """\n\nHuman: 다음의 참고자료(<context>)를 참조하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            #prompt_template = """\n\nHuman: 참고자료로 부터 구체적인 세부 정보를 충분히 제공합니다. 참고자료는 <context></context> XML tags안에 있습니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            prompt_template = """\n\nHuman: 다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>

            <question>            
            {question}
            </question>

            Assistant:"""
                
        elif conv_type == "translation":  # for translation, input
            prompt_template = """\n\nHuman: 다음의 <article>를 English로 번역하세요. 머리말은 건너뛰고 본론으로 바로 들어가주세요. 또한 결과는 <result> tag를 붙여주세요.

            <article>
            {input}
            </article>
                        
            Assistant:"""

        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: 아래의 <example> review와 Extracted Topic and sentiment 인 <result>가 있습니다.
            <example>
            객실은 작지만 깨끗하고 편안합니다. 프론트 데스크는 정말 분주했고 체크인 줄도 길었지만, 직원들은 프로페셔널하고 매우 유쾌하게 각 사람을 응대했습니다. 우리는 다시 거기에 머물것입니다.
            </example>
            <result>
            청소: 긍정적, 
            서비스: 긍정적
            </result>

            아래의 <review>에 대해서 위의 <result> 예시처럼 Extracted Topic and sentiment 을 만들어 주세요..

            <review>
            {input}
            </review>

            Assistant:"""

        elif conv_type == "extraction":  # information extraction
            prompt_template = """\n\nHuman: 다음 텍스트에서 이메일 주소를 정확하게 복사하여 한 줄에 하나씩 적어주세요. 입력 텍스트에 정확하게 쓰여있는 이메일 주소만 적어주세요. 텍스트에 이메일 주소가 없다면, "N/A"라고 적어주세요. 또한 결과는 <result> tag를 붙여주세요.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: 아래의 <text>에서 개인식별정보(PII)를 모두 제거하여 외부 계약자와 안전하게 공유할 수 있도록 합니다. 이름, 전화번호, 주소, 이메일을 XXX로 대체합니다. 또한 결과는 <result> tag를 붙여주세요.
            
            <text>
            {input}
            </text>
        
            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: 다음의 <article>에서 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요.

            <article>
            {input}
            </article>
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. 아래 문맥(context)을 참조했음에도 답을 알 수 없다면, 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다.

            {input}

            Assistant: 단계별로 생각할까요?

            Human: 예, 그렇게하세요.
            
            Assistant:"""

        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: 다음 대화를 완성하기 위해 "A"로 말하는 다음 줄을 작성하세요. Assistant는 유치원 선생님처럼 대화를 합니다.
            
            Q: 이빨 요정은 실제로 있을까?
            A: 물론이죠, 오늘 밤 당신의 이를 감싸서 베개 밑에 넣어두세요. 아침에 뭔가 당신을 기다리고 있을지도 모릅니다.
            Q: {input}

            Assistant:"""      

        elif conv_type == "timestamp-extraction":
            prompt_template = """\n\nHuman: 아래의 <text>는 시간을 포함한 텍스트입니다. 친절한 AI Assistant로서 시간을 추출하여 아래를 참조하여 <example>과 같이 정리해주세요.
            
            - 년도를 추출해서 <year>/<year>로 넣을것 
            - 월을 추출해서 <month>/<month>로 넣을것
            - 일을 추출해서 <day>/<day>로 넣을것
            - 시간을 추출해서 24H으로 정리해서 <hour>/<hour>에 넣을것
            - 분을 추출해서 <minute>/<minute>로 넣을것

            이때의 예제는 아래와 같습니다.
            <example>
            2022년 11월 3일 18시 26분
            </example>
            <result>
                <year>2022</year>
                <month>11</month>
                <day>03</day>
                <hour>18</hour>
                <minute>26</minute>
            </result>

            결과에 개행문자인 "\n"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.

            <text>
            {input}
            </text>

            Assistant:"""  

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. 모든 대화는 반말로하여야 합니다. Assistant의 이름은 서서이고 10살 여자 어린이 상상력이 풍부하고 재미있는 대화를 합니다. 때로는 바보같은 답변을 해서 재미있게 해줍니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        elif conv_type == "get-weather":  # getting weather (function calling)
            prompt_template = """\n\nHuman: In this environment you have access to a set of tools you can use to answer the user's question.

            You may call them like this. Only invoke one function at a time and wait for the results before invoking another function:
            
            <function_calls>
            <invoke>
            <tool_name>$TOOL_NAME</tool_name>
            <parameters>
            <$PARAMETER_NAME>$PARAMETER_VALUE</$PARAMETER_NAME>
            ...
            </parameters>
            </invoke>
            </function_calls>

            Here are the tools available:
            <tools>
            {tools_string}
            </tools>

            Human:
            {user_input}

            Assistant:"""                  
                
        else:
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant는 모르는 질문을 받으면 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다. 
        
            <question>            
            {question}
            </question>

            Assistant:"""

    else:  # English
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.

            <history>
            {history}
            </history>
            
            <question>            
            {input}
            </question>

            Assistant:"""

        elif conv_type=='qa':  # for RAG
            prompt_template = """\n\nHuman: Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer. 
            
            <context>
            {context}
            </context>

            Go directly into the main points without the preamble. Do not include any additional information like newline characters "\n" or character counts in the result.
                        
            <question>
            {question}
            </question>

            Assistant:"""

        elif conv_type=="translation": 
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags. Translate the article to Korean. Put it in <result> tags.
            
            <article>
            {input}
            </article>
                        
            Assistant:"""
        
        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: Here is <example> review and extracted topics and sentiments as <result>.

            <example>
            The room was small but clean and comfortable. The front desk was really busy and the check-in line was long, but the staff were professional and very pleasant with each person they helped. We will stay there again.
            </example>

            <result>
            Cleanliness: Positive, 
            Service: Positive
            </result>

            <review>
            {input}
            </review>
            
            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: We want to de-identify some text by removing all personally identifiable information from this text so that it can be shared safely with external contractors.
            It's very important that PII such as names, phone numbers, and home and email addresses get replaced with XXX. Put it in <result> tags.

            Here is the text, inside <text></text> XML tags.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "extraction":  # for sentiment, input
            prompt_template = """\n\nHuman: Please precisely copy any email addresses from the following text and then write them, one per line.  Only write an email address if it's precisely spelled out in the input text.  If there are no email addresses in the text, write "N/A".  Do not say anything else.  Put it in <result> tags.

            {input}

            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags:

            <article>
            {input}
            </article>

            Please identify any grammatical errors in the article. Also, add the fixed article at the end of answer.
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.
            
            {input}

            Assistant: Can I think step by step?

            Human: Yes, please do.

            Assistant:"""
        
        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: Please complete the conversation by writing the next line, speaking as "A". You will be acting as a kindergarten teacher.

            Q: Is the tooth fairy real?
            A: Of course, sweetie. Wrap up your tooth and put it under your pillow tonight. There might be something waiting for you in the morning.
            Q: {input}

            Assistant:"""       

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant의 이름은 서서이고 10살 여자 어린이입니다. 상상력이 풍부하고 재미있는 대화를 잘합니다. 때론 바보같은 답변을 합니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        else: # normal
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor named Seoyeon.

            Human: {input}

            Assistant:"""

            # Use the following pieces of context to provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer. 
            # The following is a friendly conversation between a human and an AI. The AI is talkative and provides lots of specific details from its context. If the AI does not know the answer to a question, it truthfully says it does not know.
    
    return PromptTemplate.from_template(prompt_template)

# load documents from s3 
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt' or file_type == 'md':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
                    
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load a code file from s3
def load_code(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    separators=""
    if file_type == 'py':        
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\ndef "]
        #print('contents: ', contents)
    elif file_type == 'js':
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\nfunction ", "\nexports.handler "]
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=50,
        chunk_overlap=0,
        #separators=["def ", "\n\n", "\n", ".", " ", ""],
        separators=separators,
        length_function = len,
    ) 

    texts = text_splitter.split_text(contents) 
    
    for i, text in enumerate(texts):
        print(f"Chunk #{i}: {text}")
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs

def get_summary(chat, docs):    
    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내로 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters."
        )
        #prompt_template = """\n\nHuman: 다음 텍스트를 간결하게 요약하세오. 텍스트의 요점을 다루는 글머리 기호로 응답을 반환합니다.
        #prompt_template = """\n\nHuman: 아래 <text>는 문서에서 추출한 텍스트입니다. 친절한 AI Assistant로서 아래와 같이 정리해주세요.
        
        #- 50자 미안의 제목을 <title>Name: </title> 안에 넣을것
        #- 300자 미안의 설명을 <description>설명: </description> 안에 넣을것
        #- 500자 미만의 내용 요약을 <summarization>요약: </summarization> 안에 넣을것
        #- 10자 미안의 애용과 과련된 테그 5개를 <tag></tag> 테그 안에 생성할 것

        #모든 생성 결과는 한국어로 해주세요. 결과에 개행문자인 "\m"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.
        #생성이 어렵거나 해당 내용을 모르는 경우 "None"로 결과를 생성하세요.
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        print('result of summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def generate_code(connectionId, requestId, chat, text, context, mode):
    if mode == 'py':    
        system = (
            """다음의 <context> tag안에는 질문과 관련된 python code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 python 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
            
            <context>
            {context}
            </context>"""
        )
    elif mode == 'js':
        system = (
            """다음의 <context> tag안에는 질문과 관련된 node.js code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 node.js 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "<context>{text}</context>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "context": context,
                "text": text
            }
        )
        
        geenerated_code = readStreamMsg(connectionId, requestId, stream.content)
                              
        geenerated_code = stream.content        
        print('result of code generation: ', geenerated_code)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return geenerated_code

def summary_of_code(chat, code, mode):
    if mode == 'py':
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def revise_question(connectionId, requestId, chat, query):    
    global history_length, token_counter_history    
    history_length = token_counter_history = 0
        
    if isKorean(query)==True :      
        system = (
            ""
        )  
        human = """이전 대화를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다. 결과는 <result> tag를 붙여주세요.
        
        <question>            
        {question}
        </question>"""
        
    else: 
        system = (
            ""
        )
        human = """Rephrase the follow up <question> to be a standalone question. Put it in <result> tags.
        <question>            
        {question}
        </question>"""
            
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "history": history,
                "question": query,
            }
        )
        generated_question = result.content
        
        revised_question = generated_question[generated_question.find('<result>')+8:len(generated_question)-9] # remove <result> tag                   
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
            
        sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n * 대화이력({str(history_length)}자, {token_counter_history} Tokens)을 활용하였습니다.")
            
    return revised_question    
    # return revised_question.replace("\n"," ")

def query_using_RAG_context(connectionId, requestId, chat, context, revised_question):    
    if isKorean(revised_question)==True:
        system = (
            """다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>"""
        )
    else: 
        system = (
            """Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
                   
    chain = prompt | chat
    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "context": context,
                "input": revised_question,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
        print('msg: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    return msg
    
def extract_text(chat, img_base64):    
    query = "텍스트를 추출해서 utf8로 변환하세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        extracted_text = result.content
        print('result of text extraction from an image: ', extracted_text)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return extracted_text
    
def load_chat_history(userId, allowTime):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text' and text and msg:
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            #print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

kendraRetriever = AmazonKendraRetriever(
    index_id=kendraIndex, 
    top_k=top_k, 
    region_name=kendra_region,
    attribute_filter = {
        "EqualsTo": {      
            "Key": "_language_code",
            "Value": {
                "StringValue": "ko"
            }
        },
    },
)

def retrieve_from_kendra(query, top_k):
    if kendra_method == 'kendra_retriever':
        relevant_docs = retrieve_from_kendra_using_kendra_retriever(query, top_k)
    else: 
        relevant_docs = retrieve_from_kendra_using_custom_retriever(query, top_k)
    
    return relevant_docs

def retrieve_from_kendra_using_kendra_retriever(query, top_k):
    print(f"query: {query} (kendra)")

    relevant_docs = []
    relevant_documents = kendraRetriever.get_relevant_documents(
        query=query,
        top_k=top_k,
    )
    #print('length of relevant_documents: ', len(relevant_documents))
    #print('relevant_documents: ', relevant_documents)

    rag_type = "kendra"
    api_type = "kendraRetriever"

    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        # print(f'## Document(retrieve_from_kendra_using_kendra_retriever) {i+1}: {document}')

        result_id = document.metadata['result_id']
        document_id = document.metadata['document_id']
        # source = document.metadata['source']
        title = document.metadata['title']
        excerpt = document.metadata['excerpt']

        uri = ""
        if "_source_uri" in document.metadata['document_attributes']:
            uri = document.metadata['document_attributes']['_source_uri']

        page = ""
        if "_excerpt_page_number" in document.metadata['document_attributes']:            
            page = document.metadata['document_attributes']['_excerpt_page_number']

        confidence = ""
        assessed_score = ""
            
        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    #"type": query_result_type,
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",                    
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                #"query_id": query_id,
                #"feedback_token": feedback_token
                "assessed_score": assessed_score,
                "result_id": result_id
            }

        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    #"type": query_result_type,
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                    "translated_excerpt": ""
                },
                #"query_id": query_id,
                #"feedback_token": feedback_token
                "assessed_score": assessed_score,
                "result_id": result_id
            }
            
        relevant_docs.append(doc_info)
    
    return relevant_docs    
    
def retrieve_from_kendra_using_custom_retriever(query, top_k):
    print(f"query: {query} (kendra)")

    index_id = kendraIndex    
    
    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    try:
        resp =  kendra_client.retrieve(
            IndexId = index_id,
            QueryText = query,
            PageSize = top_k,      
            AttributeFilter = {
                "EqualsTo": {      
                    "Key": "_language_code",
                    "Value": {
                        "StringValue": "ko"
                    }
                },
            },      
        )
        # print('retrieve resp:', json.dumps(resp))
        query_id = resp["QueryId"]

        if len(resp["ResultItems"]) >= 1:
            relevant_docs = []
            retrieve_docs = []
            for query_result in resp["ResultItems"]:
                #confidence = query_result["ScoreAttributes"]['ScoreConfidence']
                #if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': only for "en"
                retrieve_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="retrieve", query_result=query_result))
                # print('retrieve_docs: ', retrieve_docs)

            print('Looking for FAQ...')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = 4, # Maximum number of results returned for FAQ = 4 (default)
                    QueryResultTypeFilter = "QUESTION_ANSWER",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', json.dumps(resp))
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        #if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                        if confidence == 'VERY_HIGH' or confidence == 'HIGH': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="query", query_result=query_result))

                            if len(relevant_docs)>=top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for FAQ')

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")

            for doc in retrieve_docs:                
                if len(relevant_docs)>=top_k:
                    break
                else:
                    relevant_docs.append(doc)
            
        else:  # fallback using query API
            print('No result for Retrieve API!')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = top_k,
                    #QueryResultTypeFilter = "DOCUMENT",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', resp)
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    relevant_docs = []
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="query", query_result=query_result))

                            if len(relevant_docs)>=top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for Query API. Finally, no relevant docs!')
                    relevant_docs = []

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")                

    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to retrieve from Kendra")     

    #for i, rel_doc in enumerate(relevant_docs):
    #    print(f'## Document(retrieve_from_kendra_using_kendra_retriever) {i+1}: {json.dumps(rel_doc)}')  

    return relevant_docs

def priority_search(query, relevant_docs, minSimilarity):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        if 'translated_excerpt' in doc['metadata'] and doc['metadata']['translated_excerpt']:
            content = doc['metadata']['translated_excerpt']
        else:
            content = doc['metadata']['excerpt']
            
        # print('content: ', content)
        
        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    # print('excerpts: ', excerpts)

    embeddings = get_ps_embedding()
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        #k=top_k
        k=len(relevant_docs)
    )

    docs = []
    for i, document in enumerate(rel_documents):
        # print(f'## Document(priority_search) {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < minSimilarity:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

def extract_relevant_doc_for_kendra(query_id, api_type, query_result):
    rag_type = "kendra"
    if(api_type == 'retrieve'): # retrieve API
        excerpt = query_result["Content"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = query_result["DocumentTitle"]
        
        document_uri = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_source_uri":
                document_uri = str(attribute["Value"]["StringValue"])        
        if document_uri=="":  
            document_uri = query_result["DocumentURI"]

        doc_info = {
            "rag_type": rag_type,
            "api_type": api_type,
            "confidence": confidence,
            "metadata": {
                "document_id": document_id,
                "source": document_uri,
                "title": document_title,
                "excerpt": excerpt,
                "translated_excerpt": ""
            },
            "assessed_score": "",
        }
            
    else: # query API
        query_result_type = query_result["Type"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = ""
        if "Text" in query_result["DocumentTitle"]:
            document_title = query_result["DocumentTitle"]["Text"]
        document_uri = query_result["DocumentURI"]
        feedback_token = query_result["FeedbackToken"] 

        page = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_excerpt_page_number":
                page = str(attribute["Value"]["LongValue"])

        if query_result_type == "QUESTION_ANSWER":
            question_text = ""
            additional_attributes = query_result["AdditionalAttributes"]
            for attribute in additional_attributes:
                if attribute["Key"] == "QuestionText":
                    question_text = str(attribute["Value"]["TextWithHighlightsValue"]["Text"])
            answer = query_result["DocumentExcerpt"]["Text"]
            excerpt = f"{question_text} {answer}"
            excerpt = excerpt.replace("\n"," ") 
        else: 
            excerpt = query_result["DocumentExcerpt"]["Text"]

        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
    return doc_info

def get_reference(docs, rag_method, rag_type, path, doc_prefix):
    if rag_method == 'RetrievalQA' or rag_method == 'ConversationalRetrievalChain':
        if rag_type == 'kendra':
            reference = "\n\nFrom\n"
            for i, doc in enumerate(docs):
                name = doc.metadata['title']     

                uri = ""
                if ("document_attributes" in doc.metadata) and ("_source_uri" in doc.metadata['document_attributes']):
                    uri = doc.metadata['document_attributes']['_source_uri']
                                    
                if ("document_attributes" in doc.metadata) and ("_excerpt_page_number" in doc.metadata['document_attributes']):
                    page = doc.metadata['document_attributes']['_excerpt_page_number']
                    reference = reference + f'{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>\n'
                else:
                    reference = reference + f'{i+1}. <a href={uri} target=_blank>{name}</a>\n'
        else:
            reference = "\n\nFrom\n"
            for i, doc in enumerate(docs):
                # print(f'## Document(get_reference) {i+1}: {doc}')

                name = doc.metadata['name']
                page = doc.metadata['page']
                uri = doc.metadata['uri']

                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>\n"

    elif rag_method == 'RetrievalPrompt':
        reference = "\n\nFrom\n"
        for i, doc in enumerate(docs):
            if doc['metadata']['translated_excerpt']:
                excerpt = str(doc['metadata']['excerpt']+'  [번역]'+doc['metadata']['translated_excerpt']).replace('"',"") 
            else:
                excerpt = str(doc['metadata']['excerpt']).replace('"'," ")
            
            excerpt = excerpt.replace('\n','\\n')                        
                
            if doc['rag_type'] == 'kendra':                
                if doc['api_type'] == 'kendraRetriever': # provided by kendraRetriever from langchain
                    name = doc['metadata']['title']
                    uri = doc['metadata']['source']
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                elif doc['api_type'] == 'retrieve': # Retrieve. socre of confidence is only avaialbe for English
                    uri = doc['metadata']['source']
                    name = doc['metadata']['title']
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                else: # Query
                    confidence = doc['confidence']
                    if ("type" in doc['metadata']) and (doc['metadata']['type'] == "QUESTION_ANSWER"):
                        reference = reference + f"{i+1}. <a href=\"#\" onClick=\"alert(`{excerpt}`)\">FAQ ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    else:
                        uri = ""
                        if "title" in doc['metadata']:
                            #print('metadata: ', json.dumps(doc['metadata']))
                            name = doc['metadata']['title']
                            if name: 
                                uri = path+doc_prefix+parse.quote(name)

                        page = ""
                        if "document_attributes" in doc['metadata']:
                            if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                                page = doc['metadata']['document_attributes']['_excerpt_page_number']
                                                
                        if page: 
                            reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name} ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                        elif uri:
                            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name} ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
            elif doc['rag_type'][:10] == 'opensearch':
                # print(f'## Document(get_reference) {i+1}: {doc}')
                
                page = ""
                if "document_attributes" in doc['metadata']:
                    if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                        page = doc['metadata']['document_attributes']['_excerpt_page_number']
                uri = doc['metadata']['source']
                name = doc['metadata']['title']

                #print('opensearch page: ', page)

                if page:                
                    reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                else:
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    
            elif doc['rag_type'] == 'search': # google search
                # print(f'## Document(get_reference) {i+1}: {doc}')
                
                uri = doc['metadata']['source']
                name = doc['metadata']['title']
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        
    return reference

def get_code_reference(docs):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        excerpt = doc['metadata']['excerpt'].replace('"','')
        code = doc['metadata']['code'].replace('"','')
        
        excerpt = excerpt.replace('\n','\\n')
        code = code.replace('\n','\\n')
        print('reference_doc: ', json.dumps(doc))
        
        if doc['rag_type'][:10] == 'opensearch':
            # print(f'## Document(get_reference) {i+1}: {doc}')
                
            page = ""
            if "document_attributes" in doc['metadata']:
                if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                    page = doc['metadata']['document_attributes']['_excerpt_page_number']
            uri = doc['metadata']['source']
            name = doc['metadata']['title']
            name = name[name.rfind('/')+1:len(name)]

            if page:                
                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">코드설명</a>, <a href=\"#\" onClick=\"alert(`{code}`)\">관련코드</a>\n"
            else:
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">코드설명</a>, <a href=\"#\" onClick=\"alert(`{code}`)\">관련코드</a>\n"
                            
    return reference

def get_parent_content(parent_doc_id):
    response = os_client.get(
        index="idx-rag", 
        id = parent_doc_id
    )
    
    source = response['_source']                            
    # print('parent_doc: ', source['text'])   
    
    metadata = source['metadata']    
    #print('name: ', metadata['name'])   
    #print('uri: ', metadata['uri'])   
    #print('doc_level: ', metadata['doc_level']) 
    
    return source['text'], metadata['name'], metadata['uri']

def get_parent_document(doc):
    # print('doc: ', doc)
    if 'parent_doc_id' in doc['metadata']:
        parent_doc_id = doc['metadata']['parent_doc_id']
    
        if parent_doc_id:
            response = os_client.get(
                index="idx-rag", 
                id = parent_doc_id
            )
            
            #source = response['_source']
            # print('parent_doc: ', source['text'])   
            
            #metadata = source['metadata']    
            #print('name: ', metadata['name'])   
            #print('uri: ', metadata['uri'])   
            #print('doc_level: ', metadata['doc_level']) 
            
            print('text(before)', doc['metadata']['excerpt'])
            doc['metadata']['excerpt'] = response['_source']['text']
            print('text(after)', doc['metadata']['excerpt'])
        
    return doc
    
def retrieve_docs_from_vectorstore(vectorstore_opensearch, query, top_k, rag_type):
    print(f"query: {query} ({rag_type})")

    rel_docs_vector_search = []
    rel_docs_lexical_search = []        
    if rag_type == 'opensearch':                                                        
        # vector search (semantic) 
        if enalbeParentDocumentRetrival=='true':
            result = vectorstore_opensearch.similarity_search_with_score(
                query = query,
                k = top_k*2,  # use double
                pre_filter={"doc_level": {"$eq": "child"}}
            )
            print('result of opensearch: ', result)
                    
            relevant_documents = []
            docList = []
            for re in result:
                if 'parent_doc_id' in re[0].metadata:
                    parent_doc_id = re[0].metadata['parent_doc_id']
                    doc_level = re[0].metadata['doc_level']
                    print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                            
                    if doc_level == 'child':
                        if parent_doc_id in docList:
                            print('duplicated!')
                        else:
                            relevant_documents.append(re)
                            docList.append(parent_doc_id)
                            
                            #if len(relevant_documents)>=top_k:
                            #    break    
        else:
            relevant_documents = vectorstore_opensearch.similarity_search_with_score(
                query = query,
                k = top_k,
            )
        #print('(opensearch score) relevant_documents: ', relevant_documents)

        for i, document in enumerate(relevant_documents):
            #print('document.page_content:', document.page_content)
            #print('document.metadata:', document.metadata)
            # print(f'## Document(opensearch-vector) {i+1}: {document}')

            name = document[0].metadata['name']
            # print('metadata: ', document[0].metadata)

            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']

            excerpt = document[0].page_content
            confidence = str(document[1])
            assessed_score = str(document[1])
            
            parent_doc_id = doc_level = ""            
            if enalbeParentDocumentRetrival == 'true':
                parent_doc_id = document[0].metadata['parent_doc_id']
                doc_level = document[0].metadata['doc_level']
                
            if page:
                print('page: ', page)
                doc_info = {
                    "rag_type": 'opensearch-vector',
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "document_attributes": {
                            "_excerpt_page_number": page
                        },
                        "parent_doc_id": parent_doc_id,
                        "doc_level": doc_level                        
                    },
                    "assessed_score": assessed_score,
                }
            else:
                doc_info = {
                    "rag_type": 'opensearch-vector',
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "parent_doc_id": parent_doc_id,
                        "doc_level": doc_level
                    },
                    "assessed_score": assessed_score,
                }
            rel_docs_vector_search.append(doc_info)
        print(f'rel_docs (vector): '+json.dumps(rel_docs_vector_search))
    
        # lexical search (keyword)
        min_match = 0
        if enableHybridSearch == 'true':
            query = {
                "query": {
                    "bool": {
                        "must": [
                            {
                                "match": {
                                    "text": {
                                        "query": query,
                                        "minimum_should_match": f'{min_match}%',
                                        "operator":  "or",
                                        # "fuzziness": "AUTO",
                                        # "fuzzy_transpositions": True,
                                        # "zero_terms_query": "none",
                                        # "lenient": False,
                                        # "prefix_length": 0,
                                        # "max_expansions": 50,
                                        # "boost": 1
                                    }
                                }
                            },
                        ],
                        "filter": [
                        ]
                    }
                }
            }

            response = os_client.search(
                body=query,
                index="idx-*", # all
            )
            # print('lexical query result: ', json.dumps(response))
            
            for i, document in enumerate(response['hits']['hits']):
                if i>=top_k: 
                    break
                
                excerpt = document['_source']['text']
                # print(f'## Document(opensearch-keyword) {i+1}: {excerpt}')

                name = document['_source']['metadata']['name']
                # print('name: ', name)

                page = ""
                if "page" in document['_source']['metadata']:
                    page = document['_source']['metadata']['page']
                
                uri = ""
                if "uri" in document['_source']['metadata']:
                    uri = document['_source']['metadata']['uri']
                # print('uri: ', uri)

                confidence = str(document['_score'])
                assessed_score = ""
                
                parent_doc_id = doc_level = ""            
                if enalbeParentDocumentRetrival == 'true':
                    if 'parent_doc_id' in document['_source']['metadata']:
                        parent_doc_id = document['_source']['metadata']['parent_doc_id']
                    if 'doc_level' in document['_source']['metadata']:
                        doc_level = document['_source']['metadata']['doc_level']
                    
                if page:
                    print('page: ', page)
                    doc_info = {
                        "rag_type": 'opensearch-keyword',
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": name,
                            "excerpt": excerpt,
                            "translated_excerpt": "",
                            "document_attributes": {
                                "_excerpt_page_number": page
                            },
                            "parent_doc_id": parent_doc_id,
                            "doc_level": doc_level
                        },
                        "assessed_score": assessed_score,
                    }
                else: 
                    doc_info = {
                        "rag_type": 'opensearch-keyword',
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": name,
                            "excerpt": excerpt,
                            "translated_excerpt": "",
                            "parent_doc_id": parent_doc_id,
                            "doc_level": doc_level
                        },
                        "assessed_score": assessed_score,
                    }
                rel_docs_lexical_search.append(doc_info)
            print(f'rel_docs (lexical): '+json.dumps(rel_docs_lexical_search))
    
    relevant_docs = rel_docs_vector_search + rel_docs_lexical_search

    return relevant_docs

def checkDupulication(relevant_codes, doc_info):
    for doc in relevant_codes:
        if doc['metadata']['excerpt'] == doc_info['metadata']['excerpt']:
            return True
    return False

def retrieve_codes_from_vectorstore(vectorstore_opensearch, index_name, query, top_k, rag_type):
    print(f"query: {query} ({rag_type})")
    relevant_codes = []
        
    if rag_type == 'opensearch':
        # Vector Search
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        #print('(opensearch score) relevant_documents: ', relevant_documents)

        for i, document in enumerate(relevant_documents):
            # print(f'## Document(opensearch-vector) {i+1}: {document}')

            if "code" in document[0].metadata:
                code = document[0].metadata['code']
                
                name = document[0].metadata['name']
                # print('metadata: ', document[0].metadata)

                page = ""
                if "page" in document[0].metadata:
                    page = document[0].metadata['page']
                uri = ""
                if "uri" in document[0].metadata:
                    uri = document[0].metadata['uri']

                excerpt = document[0].page_content
                confidence = str(document[1])
                assessed_score = str(document[1])
                                    
                function_name = ""
                if "function_name" in document[0].metadata:
                    function_name = document[0].metadata['function_name']

                if page:
                    print('page: ', page)
                    doc_info = {
                        "rag_type": 'opensearch-vector',
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": name,
                            "excerpt": excerpt,
                            "document_attributes": {
                                "_excerpt_page_number": page
                            },
                            "code": code,
                            "function_name": function_name
                        },
                        "assessed_score": assessed_score,
                    }
                else:
                    doc_info = {
                        "rag_type": 'opensearch-vector',
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": name,
                            "excerpt": excerpt,
                            "code": code,
                            "function_name": function_name
                        },
                        "assessed_score": assessed_score,
                    }
                relevant_codes.append(doc_info)
            else:
                print("No code in metadata")
    
        # Lexical Search (keyword)
        min_match = 0
        if enableHybridSearch == 'true':
            query = {
                "query": {
                    "bool": {
                        "must": [
                            {
                                "match": {
                                    "text": {
                                        "query": query,
                                        "minimum_should_match": f'{min_match}%',
                                        "operator":  "or"
                                    }
                                }
                            },
                        ],
                        "filter": [
                        ]
                    }
                }
            }

            response = os_client.search(
                body=query,
                index=index_name, # all
            )
            # print('lexical query result: ', json.dumps(response))
            
            for i, document in enumerate(response['hits']['hits']):
                if i>=top_k: 
                    break
                
                if "code" in document['_source']['metadata']:          
                    code = document['_source']['metadata']['code']
                              
                    excerpt = document['_source']['text']
                    # print(f'## Document(opensearch-keyward) {i+1}: {excerpt}')

                    name = document['_source']['metadata']['name']
                    # print('name: ', name)

                    page = ""
                    if "page" in document['_source']['metadata']:
                        page = document['_source']['metadata']['page']
                    
                    uri = ""
                    if "uri" in document['_source']['metadata']:
                        uri = document['_source']['metadata']['uri']
                    # print('uri: ', uri)

                    confidence = str(document['_score'])
                    assessed_score = ""
                    
                    function_name = ""
                    if "function_name" in document['_source']['metadata']:
                        function_name = document['_source']['metadata']['function_name']

                    if page:
                        print('page: ', page)
                        doc_info = {
                            "rag_type": 'opensearch-keyward',
                            "confidence": confidence,
                            "metadata": {
                                "source": uri,
                                "title": name,
                                "excerpt": excerpt,
                                "document_attributes": {
                                    "_excerpt_page_number": page
                                },
                                "code": code,
                                "function_name": function_name
                            },
                            "assessed_score": assessed_score,
                        }
                    else: 
                        doc_info = {
                            "rag_type": 'opensearch-keyward',
                            "confidence": confidence,
                            "metadata": {
                                "source": uri,
                                "title": name,
                                "excerpt": excerpt,
                                "code": code,
                                "function_name": function_name
                            },
                            "assessed_score": assessed_score,
                        }
                    
                    if checkDupulication(relevant_codes, doc_info) == False:
                        relevant_codes.append(doc_info)
                else:
                    print("No code in metadata")
                    
    return relevant_codes

def retrieve_process_from_RAG(conn, vectorstore_opensearch, query, top_k, rag_type):
    relevant_docs = []
    if rag_type == 'kendra':
        rel_docs = retrieve_from_kendra(query=query, top_k=top_k)      
        print('rel_docs (kendra): '+json.dumps(rel_docs))
    else:
        rel_docs = retrieve_docs_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, query=query, top_k=top_k, rag_type=rag_type)
        print(f'rel_docs ({rag_type}): '+json.dumps(rel_docs))

    if(len(rel_docs)>=1):
        for doc in rel_docs:
            relevant_docs.append(doc)  
    
    conn.send(relevant_docs)
    conn.close()

def get_relevant_documents_using_parallel_processing(vectorstore_opensearch, question, top_k):
    relevant_docs = []    

    processes = []
    parent_connections = []
    for rag in capabilities:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        process = Process(target=retrieve_process_from_RAG, args=(child_conn, vectorstore_opensearch, question, top_k, rag))
        processes.append(process)

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        rel_docs = parent_conn.recv()

        if(len(rel_docs)>=1):
            for doc in rel_docs:
                relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def translate_process_from_relevent_doc(conn, chat, doc, bedrock_region):
    try: 
        translated_excerpt = traslation_to_korean(chat=chat, text=doc['metadata']['excerpt'])
        print(f"translated_excerpt ({bedrock_region}): {translated_excerpt}")

        doc['metadata']['translated_excerpt'] = translated_excerpt
    
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)       
        raise Exception (f"Not able to translate: {doc}")   
    
    conn.send(doc)
    conn.close()

def translate_relevant_documents_using_parallel_processing(docs):
    relevant_docs = []    
    processes = []
    parent_connections = []
    for doc in docs:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        chat = get_chat()
        bedrock_region = LLM_for_chat[selected_chat]['bedrock_region']

        process = Process(target=translate_process_from_relevent_doc, args=(child_conn, chat, doc, bedrock_region))
        processes.append(process)

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def get_reference_of_knoweledge_base(docs, path, doc_prefix):
    reference = "\n\nFrom\n"
    #print('path: ', path)
    #print('doc_prefix: ', doc_prefix)
    #print('prefix: ', f"/{doc_prefix}")
    
    for i, document in enumerate(docs):
        if document.page_content:
            excerpt = document.page_content
        
        score = document.metadata["score"]
        #print('score:', score)
            
        uri = document.metadata["location"]["s3Location"]["uri"] if document.metadata["location"]["s3Location"]["uri"] is not None else ""
        #print('uri:', uri)
        
        pos = uri.find(f"/{doc_prefix}")
        name = uri[pos+len(doc_prefix)+1:]
        encoded_name = parse.quote(name)
        #print('name:', name)
        
        uri = f"{path}{doc_prefix}{encoded_name}"
        #print('uri:', uri)
        
        reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    
    return reference

def run_bedrock_agent(text, connectionId, requestId):
    client_runtime = boto3.client('bedrock-agent-runtime')
    response =  client_runtime.invoke_agent(
        agentAliasId='CEXQFZT1EL',
        agentId='2SI1ONTVMW',
        inputText=text,
        sessionId='session-01',
        # memoryId='memory-01'
    )
    print('response of invoke_agent(): ', response)
    
    response_stream = response['completion']
    
    msg = ""
    try:
        for event in response_stream:
            chunk = event.get('chunk')
            if chunk:
                msg += chunk.get('bytes').decode()
                print('event: ', chunk.get('bytes').decode())
                
                result = {
                    'request_id': requestId,
                    'msg': msg,
                    'status': 'proceeding'
                }
                #print('result: ', json.dumps(result))
                sendMessage(connectionId, result)
                                
    except Exception as e:
        raise Exception("unexpected event.",e)
        
    return msg

def run_RAG_prompt_flow(text, connectionId, requestId):
    print('rag_flow_id: ', rag_flow_id)
    print('rag_flow_alias: ', rag_flow_alias)
    
    client = boto3.client(service_name='bedrock-agent')   
    
    # get flow alias arn
    response_flow_aliases = client.list_flow_aliases(
        flowIdentifier=rag_flow_id
    )
    print('response_flow_aliases: ', response_flow_aliases)
    flowAliasIdentifier = ""
    flowAlias = response_flow_aliases["flowAliasSummaries"]
    for alias in flowAlias:
        print('alias: ', alias)
        if alias['name'] == rag_flow_alias:
            flowAliasIdentifier = alias['arn']
            print('flowAliasIdentifier: ', flowAliasIdentifier)
            break
    
    # invoke_flow
    client_runtime = boto3.client('bedrock-agent-runtime')
    response = client_runtime.invoke_flow(
        flowIdentifier=rag_flow_id,
        flowAliasIdentifier=flowAliasIdentifier,
        inputs=[
            {
                "content": {
                    "document": text,
                },
                "nodeName": "FlowInputNode",
                "nodeOutputName": "document"
            }
        ]
    )
    print('response of invoke_flow(): ', response)
    
    response_stream = response['responseStream']
    try:
        result = {}
        for event in response_stream:
            print('event: ', event)
            result.update(event)
        print('result: ', result)

        if result['flowCompletionEvent']['completionReason'] == 'SUCCESS':
            print("Prompt flow invocation was successful! The output of the prompt flow is as follows:\n")
            # msg = result['flowOutputEvent']['content']['document']
            
            msg = readStreamMsg(connectionId, requestId, result['flowOutputEvent']['content']['document'])
            print('msg: ', msg)
        else:
            print("The prompt flow invocation completed because of the following reason:", result['flowCompletionEvent']['completionReason'])
    except Exception as e:
        raise Exception("unexpected event.",e)

    return msg

def run_prompt_flow(text, connectionId, requestId):    
    print('flow_id: ', flow_id)
    print('flow_alias: ', flow_alias)
    
    client = boto3.client(service_name='bedrock-agent')   
    
    """
    # flow (debug)      
    response_flow = client.get_flow(
        flowIdentifier=flow_id
    )
    print('response_flow: ', response_flow)
    
    definition = response_flow['definition']
    print('definition: ', definition)
    connections = definition['connections']
    print('connections: ', connections)
    for c in connections:
        print('connection: ', c)
    """
    
    # get flow alias arn
    response_flow_aliases = client.list_flow_aliases(
        flowIdentifier=flow_id
    )
    print('response_flow_aliases: ', response_flow_aliases)
    flowAliasIdentifier = ""
    flowAlias = response_flow_aliases["flowAliasSummaries"]
    for alias in flowAlias:
        print('alias: ', alias)
        if alias['name'] == flow_alias:
            flowAliasIdentifier = alias['arn']
            print('flowAliasIdentifier: ', flowAliasIdentifier)
            break
    
    # invoke_flow
    client_runtime = boto3.client('bedrock-agent-runtime')
    response = client_runtime.invoke_flow(
        flowIdentifier=flow_id,
        flowAliasIdentifier=flowAliasIdentifier,
        inputs=[
            {
                "content": {
                    "document": text,
                },
                "nodeName": "FlowInputNode",
                "nodeOutputName": "document"
            }
        ]
    )
    print('response of invoke_flow(): ', response)
    
    response_stream = response['responseStream']
    try:
        result = {}
        for event in response_stream:
            print('event: ', event)
            result.update(event)
        print('result: ', result)

        if result['flowCompletionEvent']['completionReason'] == 'SUCCESS':
            print("Prompt flow invocation was successful! The output of the prompt flow is as follows:\n")
            # msg = result['flowOutputEvent']['content']['document']
            
            msg = readStreamMsg(connectionId, requestId, result['flowOutputEvent']['content']['document'])
            print('msg: ', msg)
        else:
            print("The prompt flow invocation completed because of the following reason:", result['flowCompletionEvent']['completionReason'])
    except Exception as e:
        raise Exception("unexpected event.",e)

    return msg

def get_answer_using_knowledge_base(chat, text, connectionId, requestId):    
    revised_question = text # use original question for test
    
    retriever = AmazonKnowledgeBasesRetriever(
        knowledge_base_id="CFVYNN0NQN", 
        retrieval_config={"vectorSearchConfiguration": {"numberOfResults": 4}},
    )
    
    relevant_docs = retriever.invoke(revised_question)
    print(relevant_docs)
    
    #selected_relevant_docs = []
    #if len(relevant_docs)>=1:
    #    print('start priority search')
    #    selected_relevant_docs = priority_search(revised_question, relevant_docs, minDocSimilarity)
    #    print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
    
    relevant_context = ""
    for i, document in enumerate(relevant_docs):
        print(f"{i}: {document}")
        if document.page_content:
            content = document.page_content
        print('score:', document.metadata["score"])
        
        uri = document.metadata["location"]["s3Location"]["uri"] if document.metadata["location"]["s3Location"]["uri"] is not None else ""
        print('uri:', uri)
        
        relevant_context = relevant_context + content + "\n\n"
    
    print('relevant_context: ', relevant_context)

    msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

    reference = get_reference_of_knoweledge_base(relevant_docs, path, doc_prefix)  
    
    return msg, reference

def get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding, rag_type):
    global time_for_revise, time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_docs  # for debug
    time_for_revise = time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_docs = 0

    global time_for_rag_inference, time_for_rag_question_translation, time_for_rag_2nd_inference, time_for_rag_translation
    time_for_rag_inference = time_for_rag_question_translation = time_for_rag_2nd_inference = time_for_rag_translation = 0
    
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = "idx-*", # all
        #index_name=f"idx-{userId}',
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    ) 
    
    reference = ""
    if rag_type == 'all': # kendra, opensearch
        start_time_for_revise = time.time()

        # revise question
        revised_question = revise_question(connectionId, requestId, chat, text)     
        print('revised_question: ', revised_question)  
        revised_question = revised_question.replace('\n', '')
        
        end_time_for_revise = time.time()
        time_for_revise = end_time_for_revise - start_time_for_revise
        print('processing time for revised question: ', time_for_revise)

        relevant_docs = [] 
        if useParallelRAG == 'true':  # parallel processing
            print('start RAG for revised question')
            relevant_docs = get_relevant_documents_using_parallel_processing(vectorstore_opensearch=vectorstore_opensearch, question=revised_question, top_k=top_k)

            end_time_for_rag_inference = time.time()
            time_for_rag_inference = end_time_for_rag_inference - end_time_for_revise
            print('processing time for RAG (Inference): ', time_for_rag_inference)

            if allowDualSearch=='true' and isKorean(text)==True:
                print('start RAG for translated revised question')
                translated_revised_question = traslation_to_english(chat=chat, text=revised_question)
                print('translated_revised_question: ', translated_revised_question)

                if debugMessageMode=='true':
                    sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n번역된 새로운 질문: {translated_revised_question}")

                end_time_for_rag_question_translation = time.time()
                time_for_rag_question_translation = end_time_for_rag_question_translation - end_time_for_rag_inference
                print('processing time for RAG (Question Translation): ', time_for_rag_question_translation)

                if allowDualSearchWithMulipleProcessing == True:
                    relevant_docs_using_translated_question = get_relevant_documents_using_parallel_processing(vectorstore_opensearch=vectorstore_opensearch, question=translated_revised_question, top_k=4)

                    end_time_for_rag_2nd_inference = time.time()
                    time_for_rag_2nd_inference = end_time_for_rag_2nd_inference - end_time_for_rag_question_translation
                    print('processing time for RAG (2nd Inference): ', time_for_rag_2nd_inference)
                    
                    docs_translation_required = []
                    if len(relevant_docs_using_translated_question)>=1:
                        for i, doc in enumerate(relevant_docs_using_translated_question):
                            if isKorean(doc)==False:
                                docs_translation_required.append(doc)
                            else:
                                print(f"original {i}: {doc}")
                                relevant_docs.append(doc)
                                           
                        translated_docs = translate_relevant_documents_using_parallel_processing(docs_translation_required)
                        for i, doc in enumerate(translated_docs):
                            print(f"#### {i} (ENG): {doc['metadata']['excerpt']}")
                            print(f"#### {i} (KOR): {doc['metadata']['translated_excerpt']}")
                            relevant_docs.append(doc)
                        
                        end_time_for_rag_translation = time.time()
                        time_for_rag_translation = end_time_for_rag_translation - end_time_for_rag_2nd_inference
                        print('processing time for RAG (translation): ', time_for_rag_translation)

                else:
                    relevant_docs_using_translated_question = []
                    for reg in capabilities:
                        if reg == 'kendra':
                            rel_docs = retrieve_from_kendra(query=translated_revised_question, top_k=top_k)
                            print('rel_docs (kendra): '+json.dumps(rel_docs))
                        else:
                            rel_docs = retrieve_docs_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, query=translated_revised_question, top_k=top_k, rag_type=reg)
                            print(f'rel_docs ({reg}): '+json.dumps(rel_docs))
                    
                        if(len(rel_docs)>=1):
                            for doc in rel_docs:
                                relevant_docs_using_translated_question.append(doc)

                    if len(relevant_docs_using_translated_question)>=1:
                        for i, doc in enumerate(relevant_docs_using_translated_question):
                            if isKorean(doc)==False:
                                translated_excerpt = traslation_to_korean(chat=chat, text=doc['metadata']['excerpt'])
                                print(f"#### {i} (ENG): {doc['metadata']['excerpt']}")
                                print(f"#### {i} (KOR): {translated_excerpt}")

                                #doc['metadata']['excerpt'] = translated_excerpt
                                doc['metadata']['translated_excerpt'] = translated_excerpt
                                relevant_docs.append(doc)
                            else:
                                print(f"original {i}: {doc}")
                                relevant_docs.append(doc)
        else: # sequencial processing
            print('start the sequencial processing for multiple RAG')
            for reg in capabilities:            
                if reg == 'kendra':
                    rel_docs = retrieve_from_kendra(query=revised_question, top_k=top_k)      
                    print('rel_docs (kendra): '+json.dumps(rel_docs))
                else:
                    rel_docs = retrieve_docs_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, query=revised_question, top_k=top_k, rag_type=reg)
                    print(f'rel_docs ({reg}): '+json.dumps(rel_docs))
                
                if(len(rel_docs)>=1):
                    for doc in rel_docs:
                        relevant_docs.append(doc)

        if debugMessageMode=='true':
            for i, doc in enumerate(relevant_docs):
                print(f"#### relevant_docs ({i}): {json.dumps(doc)}")

        end_time_for_rag = time.time()
        time_for_rag = end_time_for_rag - end_time_for_revise
        print('processing time for RAG: ', time_for_rag)

        selected_relevant_docs = []
        if len(relevant_docs)>=1:
            print('start priority search')
            selected_relevant_docs = priority_search(revised_question, relevant_docs, minDocSimilarity)
            print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))

        if len(selected_relevant_docs)==0:
            print('No relevant document! So use google api')
            api_key = google_api_key
            cse_id = google_cse_id 
            
            relevant_docs = []
            try: 
                service = build("customsearch", "v1", developerKey=api_key)
                result = service.cse().list(q=revised_question, cx=cse_id).execute()
                # print('google search result: ', result)

                if "items" in result:
                    for item in result['items']:
                        api_type = "google api"
                        excerpt = item['snippet']
                        uri = item['link']
                        title = item['title']
                        confidence = ""
                        assessed_score = ""
                        
                        doc_info = {
                            "rag_type": 'search',
                            "api_type": api_type,
                            "confidence": confidence,
                            "metadata": {
                                "source": uri,
                                "title": title,
                                "excerpt": excerpt,
                                "translated_excerpt": "",
                            },
                            "assessed_score": assessed_score,
                        }
                        relevant_docs.append(doc_info)                
            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)       

                sendErrorMessage(connectionId, requestId, "Not able to use Google API. Check the credentials")    
                #sendErrorMessage(connectionId, requestId, err_msg)    
                #raise Exception ("Not able to search using google api") 
            
            if len(relevant_docs)>=1:
                selected_relevant_docs = priority_search(revised_question, relevant_docs, minDocSimilarity)
                print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
            # print('selected_relevant_docs (google): ', selected_relevant_docs)
            
        # update doc using parent
        contentList = []
        update_docs = []
        for doc in selected_relevant_docs:
            doc = get_parent_document(doc) # use pareant document
            
            # print('excerpt: ', doc['metadata']['excerpt'])
            if doc['metadata']['excerpt'] in contentList:
                print('duplicated!')
                continue
            contentList.append(doc['metadata']['excerpt'])
            update_docs.append(doc)
            
            if len(update_docs)>=top_k:
                break
        
        print('update_docs:', json.dumps(update_docs))

        end_time_for_priority_search = time.time() 
        time_for_priority_search = end_time_for_priority_search - end_time_for_rag
        print('processing time for priority search: ', time_for_priority_search)
        number_of_relevant_docs = len(update_docs)

        relevant_context = ""
        for document in update_docs:
            if document['metadata']['translated_excerpt']:
                content = document['metadata']['translated_excerpt']
            else:
                content = document['metadata']['excerpt']
        
            relevant_context = relevant_context + content + "\n\n"
        print('relevant_context: ', relevant_context)

        # query using RAG context
        msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

        reference = ""
        if len(update_docs)>=1 and enableReference=='true':
            reference = get_reference(update_docs, rag_method, rag_type, path, doc_prefix)  

        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - end_time_for_priority_search
        print('processing time for inference: ', time_for_inference)
        
    else:        
        start_time_for_revise = time.time()

        # revise question
        revised_question = revise_question(connectionId, requestId, chat, text)     
        print('revised_question: ', revised_question) 
            
        end_time_for_revise = time.time()
        time_for_revise = end_time_for_revise - start_time_for_revise
        print('processing time for revised question: ', time_for_revise)

        if rag_type == 'kendra':
            relevant_docs = retrieve_from_kendra(query=revised_question, top_k=top_k)            
        else:
            relevant_docs = retrieve_docs_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, query=revised_question, top_k=top_k, rag_type=rag_type)
        print('relevant_docs: ', json.dumps(relevant_docs))
        
        if len(relevant_docs) >= 1:
            selected_relevant_docs = priority_search(revised_question, relevant_docs, minDocSimilarity)

        # update doc using parent
        contentList = []
        update_docs = []
        for doc in selected_relevant_docs:        
            doc = get_parent_document(doc) # use pareant document
            
            # print('excerpt: ', doc['metadata']['excerpt'])
            if doc['metadata']['excerpt'] in contentList:
                print('duplicated!')
                continue
            contentList.append(doc['metadata']['excerpt'])
            update_docs.append(doc)
            
            if len(update_docs)>=top_k:
                break
        
        print('update_docs:', json.dumps(update_docs))

        end_time_for_rag = time.time()
        time_for_rag = end_time_for_rag - end_time_for_revise
        print('processing time for RAG: ', time_for_rag)
        number_of_relevant_docs = len(update_docs)

        relevant_context = ""
        for document in update_docs:
            if document['metadata']['translated_excerpt']:
                content = document['metadata']['translated_excerpt']
            else:
                content = document['metadata']['excerpt']

            relevant_context = relevant_context + content + "\n\n"
        # print('relevant_context: ', relevant_context)

        # query using RAG context
        msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

        if len(update_docs)>=1 and enableReference=='true':
            reference = get_reference(update_docs, rag_method, rag_type, path, doc_prefix)
            
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - end_time_for_rag
        print('processing time for inference: ', time_for_inference)

    global relevant_length, token_counter_relevant_docs
    
    if debugMessageMode=='true':   # extract chat history for debug
        relevant_length = len(relevant_context)
        token_counter_relevant_docs = chat.get_num_tokens(relevant_context)

    return msg, reference

def get_code_prompt_template(code_type):    
    if code_type == 'py':    
        prompt_template = """\n\nHuman: 다음의 <context> tag안에는 질문과 관련된 python code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 python 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
                
        <context>
        {context}
        </context>

        <question>            
        {question}
        </question>

        Assistant:"""
    elif code_type == 'js':    
        prompt_template = """\n\nHuman: 다음의 <context> tag안에는 질문과 관련된 node.js code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 node.js 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
                
        <context>
        {context}
        </context>

        <question>            
        {question}
        </question>

        Assistant:"""
                    
    return PromptTemplate.from_template(prompt_template)

def get_code_using_RAG(chat, text, code_type, connectionId, requestId, bedrock_embedding):
    global time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_codes  # for debug
    time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_codes = 0
    
    category = code_type
    
    index_name =  f"idx-{category}-*"
    print('index: ', index_name)
        
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = index_name,
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    )
    
    reference = ""
    start_time_for_rag = time.time()

    rag_type = 'opensearch'
    PROMPT = get_code_prompt_template(code_type)
    print('PROMPT: ', PROMPT)        

    relevant_codes = [] 
    print('start RAG for question')
    
    relevant_codes = retrieve_codes_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, index_name=index_name, query=text, top_k=top_k, rag_type=rag_type)
    print(f'relevant_codes ({rag_type}): '+json.dumps(relevant_codes))
    
    end_time_for_rag = time.time()
    time_for_rag = end_time_for_rag - start_time_for_rag
    print('processing time for RAG: ', time_for_rag)

    selected_relevant_codes = []
    if len(relevant_codes)>=1:
        selected_relevant_codes = priority_search(text, relevant_codes, minCodeSimilarity)
        print('selected_relevant_codes: ', json.dumps(selected_relevant_codes))
    
    end_time_for_priority_search = time.time() 
    time_for_priority_search = end_time_for_priority_search - end_time_for_rag
    print('processing time for priority search: ', time_for_priority_search)
    number_of_relevant_codes = len(selected_relevant_codes)
    
    # update doc using parent
    contentList = []
    update_codes = []
    for doc in selected_relevant_codes:        
        doc = get_parent_document(doc) # use pareant document
            
        # print('excerpt: ', doc['metadata']['excerpt'])
        if doc['metadata']['excerpt'] in contentList:
            print('duplicated!')
            continue
        contentList.append(doc['metadata']['excerpt'])
        update_codes.append(doc)
            
        if len(update_codes)>=top_k:
            break
        
    print('update_docs:', json.dumps(update_codes))    

    relevant_code = ""
    for document in update_codes:
        if document['metadata']['code']:
            code = document['metadata']['code']
            relevant_code = relevant_code + code + "\n\n"            
    print('relevant_code: ', relevant_code)

    msg = generate_code(connectionId, requestId, chat, text, relevant_code, code_type)
     
    if len(selected_relevant_codes)>=1 and enableReference=='true':
        reference = get_code_reference(selected_relevant_codes)  

    end_time_for_inference = time.time()
    time_for_inference = end_time_for_inference - end_time_for_priority_search
    print('processing time for inference: ', time_for_inference)
    
    return msg, reference

def get_text_speech(path, speech_prefix, bucket, msg):
    ext = "mp3"
    polly = boto3.client('polly')
    try:
        response = polly.start_speech_synthesis_task(
            Engine='neural',
            LanguageCode='ko-KR',
            OutputFormat=ext,
            OutputS3BucketName=bucket,
            OutputS3KeyPrefix=speech_prefix,
            Text=msg,
            TextType='text',
            VoiceId='Seoyeon'        
        )
        print('response: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create voice")
    
    object = '.'+response['SynthesisTask']['TaskId']+'.'+ext
    print('object: ', object)

    return path+speech_prefix+parse.quote(object)

def traslation_to_korean(chat, text):
    input_language = "English"
    output_language = "Korean"
        
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def traslation_to_english(chat, text):
    input_language = "Korean"
    output_language = "English"
        
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def summary_of_code(chat, code, mode):
    if mode == 'py':
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

@tool 
def get_book_list(keyword: str) -> str:
    """
    Search book list by keyword and then return book list
    keyword: search keyword
    return: book list
    """
    
    keyword = keyword.replace('\'','')

    answer = ""
    url = f"https://search.kyobobook.co.kr/search?keyword={keyword}&gbCode=TOT&target=total"
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        prod_info = soup.find_all("a", attrs={"class": "prod_info"})
        
        if len(prod_info):
            answer = "추천 도서는 아래와 같습니다.\n"
            
        for prod in prod_info[:5]:
            title = prod.text.strip().replace("\n", "")       
            link = prod.get("href")
            answer = answer + f"{title}, URL: {link}\n\n"
    
    return answer
    
@tool
def get_current_time(format: str=f"%Y-%m-%d %H:%M:%S")->str:
    """Returns the current date and time in the specified format"""
    # f"%Y-%m-%d %H:%M:%S"
    
    format = format.replace('\'','')
    timestr = datetime.datetime.now(timezone('Asia/Seoul')).strftime(format)
    # print('timestr:', timestr)
    
    return timestr

@tool
def get_weather_info(city: str) -> str:
    """
    Search weather information by city name and then return weather statement.
    city: the english name of city to search
    return: weather statement
    """    
    
    city = city.replace('\n','')
    city = city.replace('\'','')
    
    chat = get_chat()
                
    if isKorean(city):
        place = traslation(chat, city, "Korean", "English")
        print('city (translated): ', place)
    else:
        place = city
        city = traslation(chat, city, "English", "Korean")
        print('city (translated): ', city)
        
    print('place: ', place)
    
    weather_str: str = f"{city}에 대한 날씨 정보가 없습니다."
    if weather_api_key:
        apiKey = weather_api_key
        lang = 'en' 
        units = 'metric' 
        api = f"https://api.openweathermap.org/data/2.5/weather?q={place}&APPID={apiKey}&lang={lang}&units={units}"
        # print('api: ', api)
                
        try:
            result = requests.get(api)
            result = json.loads(result.text)
            print('result: ', result)
        
            if 'weather' in result:
                overall = result['weather'][0]['main']
                current_temp = result['main']['temp']
                min_temp = result['main']['temp_min']
                max_temp = result['main']['temp_max']
                humidity = result['main']['humidity']
                wind_speed = result['wind']['speed']
                cloud = result['clouds']['all']
                
                weather_str = f"{city}의 현재 날씨의 특징은 {overall}이며, 현재 온도는 {current_temp}도 이고, 최저온도는 {min_temp}도, 최고 온도는 {max_temp}도 입니다. 현재 습도는 {humidity}% 이고, 바람은 초당 {wind_speed} 미터 입니다. 구름은 {cloud}% 입니다."
                #weather_str = f"Today, the overall of {city} is {overall}, current temperature is {current_temp} degree, min temperature is {min_temp} degree, highest temperature is {max_temp} degree. huminity is {humidity}%, wind status is {wind_speed} meter per second. the amount of cloud is {cloud}%."            
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                    
            # raise Exception ("Not able to request to LLM")    
        
    print('weather_str: ', weather_str)                            
    return weather_str

@tool
def search_by_tavily(keyword: str) -> str:
    """
    Search general information by keyword and then return the result as a string.
    keyword: search keyword
    return: the information of keyword
    """    
    global reference_docs
    
    answer = ""
    
    if tavily_api_key:
        keyword = keyword.replace('\'','')
        
        search = TavilySearchResults(k=3)
                    
        output = search.invoke(keyword)
        print('tavily output: ', output)
        
        for result in output:
            print('result: ', result)
            if result:
                content = result.get("content")
                url = result.get("url")
                
                reference_docs.append(
                    Document(
                        page_content=content,
                        metadata={
                            'name': 'WWW',
                            'uri': url,
                            'from': 'tavily'
                        },
                    )
                )                
                answer = answer + f"{content}, URL: {url}\n"
        
    return answer

@tool    
def search_by_opensearch(keyword: str) -> str:
    """
    Search technical information by keyword and then return the result as a string.
    keyword: search keyword
    return: the technical information of keyword
    """    
    
    print('keyword: ', keyword)
    keyword = keyword.replace('\'','')
    keyword = keyword.replace('|','')
    keyword = keyword.replace('\n','')
    print('modified keyword: ', keyword)
    
    bedrock_embedding = get_embedding()
        
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = "idx-*", # all
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    ) 
    
    answer = ""
    top_k = 2
    
    docs = [] 
    if enalbeParentDocumentRetrival == 'true': # parent/child chunking
        relevant_documents = get_documents_from_opensearch(vectorstore_opensearch, keyword, top_k)
                        
        for i, document in enumerate(relevant_documents):
            #print(f'## Document(opensearch-vector) {i+1}: {document}')
            
            parent_doc_id = document[0].metadata['parent_doc_id']
            doc_level = document[0].metadata['doc_level']
            #print(f"child: parent_doc_id: {parent_doc_id}, doc_level: {doc_level}")
            
            excerpt, name, uri = get_parent_content(parent_doc_id) # use pareant document
            #print(f"parent_doc_id: {parent_doc_id}, doc_level: {doc_level}, uri: {uri}, content: {excerpt}")
            
            docs.append(
                Document(
                    page_content=excerpt,
                    metadata={
                        'name': name,
                        'uri': uri,
                        'doc_level': doc_level,
                        'from': 'vector'
                    },
                )
            )
    else: 
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = keyword,
            k = top_k,
        )

        for i, document in enumerate(relevant_documents):
            #print(f'## Document(opensearch-vector) {i+1}: {document}')
            
            excerpt = document[0].page_content        
            uri = document[0].metadata['uri']            
            name = document[0].metadata['name']
            
            docs.append(
                Document(
                    page_content=excerpt,
                    metadata={
                        'name': name,
                        'uri': uri,
                        'from': 'vector'
                    },
                )
            )
    
    if enableHybridSearch == 'true':
        docs = docs + lexical_search_for_tool(keyword, top_k)
    
    print('doc length: ', len(docs))
                
    filtered_docs = grade_documents(keyword, docs)
        
    for i, doc in enumerate(filtered_docs):
        if len(doc.page_content)>=100:
            text = doc.page_content[:100]
        else:
            text = doc.page_content
            
        print(f"filtered doc[{i}]: {text}, metadata:{doc.metadata}")
       
    answer = "" 
    for doc in filtered_docs:
        excerpt = doc.page_content
        uri = doc.metadata['uri']
        
        answer = answer + f"{excerpt}\n\n"
        
    return answer

def get_documents_from_opensearch(vectorstore_opensearch, query, top_k):
    result = vectorstore_opensearch.similarity_search_with_score(
        query = query,
        k = top_k*2,  
        pre_filter={"doc_level": {"$eq": "child"}}
    )
    # print('result: ', result)
                
    relevant_documents = []
    docList = []
    for re in result:
        if 'parent_doc_id' in re[0].metadata:
            parent_doc_id = re[0].metadata['parent_doc_id']
            doc_level = re[0].metadata['doc_level']
            print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                        
            if doc_level == 'child':
                if parent_doc_id in docList:
                    print('duplicated!')
                else:
                    relevant_documents.append(re)
                    docList.append(parent_doc_id)
                        
                    if len(relevant_documents)>=top_k:
                        break                                
    # print('lexical query result: ', json.dumps(response))
    
    for i, doc in enumerate(relevant_documents):
        #print('doc: ', doc[0])
        #print('doc content: ', doc[0].page_content)
        
        if len(doc[0].page_content)>=30:
            text = doc[0].page_content[:30]
        else:
            text = doc[0].page_content            
        print(f"--> vector search doc[{i}]: {text}, metadata:{doc[0].metadata}")        

    return relevant_documents

def lexical_search_for_tool(query, top_k):
    # lexical search (keyword)
    min_match = 0
    
    query = {
        "query": {
            "bool": {
                "must": [
                    {
                        "match": {
                            "text": {
                                "query": query,
                                "minimum_should_match": f'{min_match}%',
                                "operator":  "or",
                            }
                        }
                    },
                ],
                "filter": [
                ]
            }
        }
    }

    response = os_client.search(
        body=query,
        index="idx-*", # all
    )
    # print('lexical query result: ', json.dumps(response))
        
    docs = []
    for i, document in enumerate(response['hits']['hits']):
        if i>=top_k: 
            break
                    
        excerpt = document['_source']['text']
        
        name = document['_source']['metadata']['name']
        # print('name: ', name)

        page = ""
        if "page" in document['_source']['metadata']:
            page = document['_source']['metadata']['page']
        
        uri = ""
        if "uri" in document['_source']['metadata']:
            uri = document['_source']['metadata']['uri']            
        
        docs.append(
                Document(
                    page_content=excerpt,
                    metadata={
                        'name': name,
                        'uri': uri,
                        'page': page,
                        'from': 'lexical'
                    },
                )
            )
    
    for i, doc in enumerate(docs):
        #print('doc: ', doc)
        #print('doc content: ', doc.page_content)
        
        if len(doc.page_content)>=30:
            text = doc.page_content[:30]
        else:
            text = doc.page_content            
        print(f"--> lexical search doc[{i}]: {text}, metadata:{doc.metadata}")   
        
    return docs

class GradeDocuments(BaseModel):
    """Binary score for relevance check on retrieved documents."""

    binary_score: str = Field(description="Documents are relevant to the question, 'yes' or 'no'")

def get_retrieval_grader(chat):
    system = """You are a grader assessing relevance of a retrieved document to a user question. \n 
    If the document contains keyword(s) or semantic meaning related to the question, grade it as relevant. \n
    Give a binary score 'yes' or 'no' score to indicate whether the document is relevant to the question."""

    grade_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system),
            ("human", "Retrieved document: \n\n {document} \n\n User question: {question}"),
        ]
    )
    
    structured_llm_grader = chat.with_structured_output(GradeDocuments)
    retrieval_grader = grade_prompt | structured_llm_grader
    return retrieval_grader

def grade_document_based_on_relevance(conn, question, doc, models, selected):     
    chat = get_multi_region_chat(models, selected)
    retrieval_grader = get_retrieval_grader(chat)       
    score = retrieval_grader.invoke({"question": question, "document": doc.page_content})
    # print(f"score: {score}")
    
    grade = score.binary_score    
    if grade == 'yes':
        print("---GRADE: DOCUMENT RELEVANT---")
        conn.send(doc)
    else:  # no
        print("---GRADE: DOCUMENT NOT RELEVANT---")
        conn.send(None)
    
    conn.close()
                                
def grade_documents_using_parallel_processing(question, documents):
    relevant_docs = []    

    processes = []
    parent_connections = []
    
    selected = 0
    models = LLM_for_chat
    for i, doc in enumerate(documents):
        #print(f"grading doc[{i}]: {doc.page_content}")        
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        process = Process(target=grade_document_based_on_relevance, args=(child_conn, question, doc, models, selected))
        processes.append(process)

        selected = selected + 1
        if selected == len(models):
            selected = 0
    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()

        if doc is not None:
            relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def grade_documents(question, documents):
    print("###### grade_documents ######")
    
    filtered_docs = []
    if useParallelRAG == 'true':  # parallel processing
        print("start grading...")
        filtered_docs = grade_documents_using_parallel_processing(question, documents)

    else:
        # Score each doc    
        chat = get_chat()
        retrieval_grader = get_retrieval_grader(chat)
        for doc in documents:
            score = retrieval_grader.invoke({"question": question, "document": doc.page_content})
            grade = score.binary_score
            # Document relevant
            if grade.lower() == "yes":
                print("---GRADE: DOCUMENT RELEVANT---")
                filtered_docs.append(doc)
            # Document not relevant
            else:
                print("---GRADE: DOCUMENT NOT RELEVANT---")
                # We do not include the document in filtered_docs
                # We set a flag to indicate that we want to run web search
                continue
    
    global reference_docs 
    reference_docs += filtered_docs    
    # print('langth of reference_docs: ', len(reference_docs))
    
    # print('len(docments): ', len(filtered_docs))    
    return filtered_docs

def get_references_for_agent(docs):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        page = ""
        if "page" in doc.metadata:
            page = doc.metadata['page']
            #print('page: ', page)            
        uri = ""
        if "uri" in doc.metadata:
            uri = doc.metadata['uri']
            #print('uri: ', uri)                
        name = ""
        if "name" in doc.metadata:
            name = doc.metadata['name']
            #print('name: ', name)     
           
        sourceType = ""
        if "from" in doc.metadata:
            sourceType = doc.metadata['from']
        else:
            sourceType = "OpenSearch"
        #print('sourceType: ', sourceType)        
        
        #if len(doc.page_content)>=1000:
        #    excerpt = ""+doc.page_content[:1000]
        #else:
        #    excerpt = ""+doc.page_content
        excerpt = ""+doc.page_content
        print('excerpt: ', excerpt)
        
        # for some of unusual case 
        #excerpt = excerpt.replace('"', '')        
        #excerpt = ''.join(c for c in excerpt if c not in '"')
        excerpt = re.sub('"', '', excerpt)
        print('excerpt(quotation removed): ', excerpt)
        
        if page:                
            reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {sourceType}, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        else:
            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {sourceType}, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
    return reference

# define tools
tools = [get_current_time, get_book_list, get_weather_info, search_by_tavily, search_by_opensearch]        

def init_enhanced_search():
    chat = get_chat() 

    model = chat.bind_tools(tools)

    class State(TypedDict):
        messages: Annotated[list, add_messages]

    tool_node = ToolNode(tools)

    def should_continue(state: State) -> Literal["continue", "end"]:
        messages = state["messages"]    
        # print('(should_continue) messages: ', messages)
            
        last_message = messages[-1]
        if not last_message.tool_calls:
            return "end"
        else:                
            return "continue"

    def call_model(state: State):
        question = state["messages"]
        print('question: ', question)
            
        if isKorean(question[0].content)==True:
            system = (
                "Assistant는 질문에 답변하기 위한 정보를 수집하는 연구원입니다."
                "Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다."
                "Assistant는 모르는 질문을 받으면 솔직히 모른다고 말합니다."
                "최종 답변에는 조사한 내용을 반드시 포함하여야 하고, <result> tag를 붙여주세요."
            )
        else: 
            system = (            
                "You are a researcher charged with providing information that can be used when making answer."
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "You will be acting as a thoughtful advisor."
                "Put it in <result> tags."
            )
                
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
        chain = prompt | model
                
        response = chain.invoke(question)
        return {"messages": [response]}

    def buildChatAgent():
        workflow = StateGraph(State)

        workflow.add_node("agent", call_model)
        workflow.add_node("action", tool_node)
            
        workflow.set_entry_point("agent")
        workflow.add_conditional_edges(
            "agent",
            should_continue,
            {
                "continue": "action",
                "end": END,
            },
        )
        workflow.add_edge("action", "agent")
        return workflow.compile()
    
    return buildChatAgent()

app_enhanced_search = init_enhanced_search()

def enhanced_search(query):
    inputs = [HumanMessage(content=query)]
    config = {"recursion_limit": 50}
        
    result = app_enhanced_search.invoke({"messages": inputs}, config)   
    print('result: ', result)
            
    message = result["messages"][-1]
    print('enhanced_search: ', message)

    return message.content[message.content.find('<result>')+8:len(message.content)-9]

####################### LangGraph #######################
# Chat Agent Executor
#########################################################
def run_agent_executor(connectionId, requestId, query):
    chatModel = get_chat() 

    model = chatModel.bind_tools(tools)

    class State(TypedDict):
        # messages: Annotated[Sequence[BaseMessage], operator.add]
        messages: Annotated[list, add_messages]

    tool_node = ToolNode(tools)

    def should_continue(state: State) -> Literal["continue", "end"]:
        print("###### should_continue ######")
        messages = state["messages"]    
        # print('(should_continue) messages: ', messages)
        
        last_message = messages[-1]
        if not last_message.tool_calls:
            return "end"
        else:                
            return "continue"

    def call_model(state: State):
        print("###### call_model ######")
        print('state: ', state["messages"])
        
        if isKorean(state["messages"][0].content)==True:
            system = (
                "당신의 이름은 서연이고, 질문에 친근한 방식으로 대답하도록 설계된 대화형 AI입니다."
                "상황에 맞는 구체적인 세부 정보를 충분히 제공합니다."
                "모르는 질문을 받으면 솔직히 모른다고 말합니다."
                "최종 답변에는 조사한 내용을 반드시 포함합니다."
            )
        else: 
            system = (            
                "You are a conversational AI designed to answer in a friendly way to a question."
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "You will be acting as a thoughtful advisor."                
            )
            
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
        chain = prompt | model
            
        response = chain.invoke(state["messages"])
        return {"messages": [response]}

    def buildChatAgent():
        workflow = StateGraph(State)

        workflow.add_node("agent", call_model)
        workflow.add_node("action", tool_node)
        workflow.add_edge(START, "agent")
        workflow.add_conditional_edges(
            "agent",
            should_continue,
            {
                "continue": "action",
                "end": END,
            },
        )
        workflow.add_edge("action", "agent")

        return workflow.compile()

    app = buildChatAgent()
        
    isTyping(connectionId, requestId)
    
    inputs = [HumanMessage(content=query)]
    config = {"recursion_limit": 50}
    
    message = ""
    for event in app.stream({"messages": inputs}, config, stream_mode="values"):   
        # print('event: ', event)
        
        message = event["messages"][-1]
        # print('message: ', message)

    msg = readStreamMsg(connectionId, requestId, message.content)

    return msg    

####################### LangGraph #######################
# Reflection Agent
#########################################################
def run_reflection_agent(connectionId, requestId, query):
    class State(TypedDict):
        # messages: Annotated[Sequence[BaseMessage], operator.add]
        messages: Annotated[list, add_messages]

    def generation(state: State):    
        prompt = ChatPromptTemplate.from_messages(
            [
                (
                    "system",
                    "당신은 5문단의 에세이 작성을 돕는 작가이고 이름은 서연입니다"
                    "사용자의 요청에 대해 최고의 에세이를 작성하세요."
                    "사용자가 에세이에 대해 평가를 하면, 이전 에세이를 수정하여 답변하세요."
                    "최종 답변에는 완성된 에세이 전체 내용을 반드시 포함하여야 하고, <result> tag를 붙여주세요.",
                ),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
        
        chat = get_chat()
        chain = prompt | chat

        response = chain.invoke(state["messages"])
        return {"messages": [response]}

    def reflection(state: State):
        messages = state["messages"]
        
        reflection_prompt = ChatPromptTemplate.from_messages(
            [
                (
                    "system",
                    "당신은 교사로서 학셍의 에세이를 평가하삽니다. 비평과 개선사항을 친절하게 설명해주세요."
                    "이때 장점, 단점, 길이, 깊이, 스타일등에 대해 충분한 정보를 제공합니다."
                    #"특히 주제에 맞는 적절한 예제가 잘 반영되어있는지 확인합니다"
                    "각 문단의 길이는 최소 200자 이상이 되도록 관련된 예제를 충분히 포함합니다.",
                ),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
        
        chat = get_chat()
        reflect = reflection_prompt | chat
        
        cls_map = {"ai": HumanMessage, "human": AIMessage}
        translated = [messages[0]] + [
            cls_map[msg.type](content=msg.content) for msg in messages[1:]
        ]
        print('translated: ', translated)
        
        res = reflect.invoke({"messages": translated})    
        response = HumanMessage(content=res.content)    
        return {"messages": [response]}

    def should_continue(state: State) -> Literal["continue", "end"]:
        messages = state["messages"]
        
        if len(messages) >= 6:   # End after 3 iterations        
            return "end"
        else:
            return "continue"

    def buildReflectionAgent():
        workflow = StateGraph(State)
        workflow.add_node("generate", generation)
        workflow.add_node("reflect", reflection)
        workflow.set_entry_point("generate")
        workflow.add_conditional_edges(
            "generate",
            should_continue,
            {
                "continue": "reflect",
                "end": END,
            },
        )

        workflow.add_edge("reflect", "generate")
        return workflow.compile()

    app = buildReflectionAgent()

    isTyping(connectionId, requestId)
    
    inputs = [HumanMessage(content=query)]
    config = {"recursion_limit": 50}
    
    msg = ""
    
    for event in app.stream({"messages": inputs}, config, stream_mode="values"):   
        print('event: ', event)
        
        message = event["messages"][-1]
        print('message: ', message)
        
        if len(event["messages"])>1:
            if msg == "":
                msg = message.content
            else:
                msg = f"{msg}\n\n{message.content}"

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)

    return msg

####################### LangGraph #######################
# Knowledge Guru
#########################################################
def run_knowledge_guru(connectionId, requestId, query):
    class State(TypedDict):
        messages: Annotated[list, add_messages]
        reflection: list
        search_queries: list
            
    def generate(state: State):    
        print("###### generate ######")
        print('state: ', state["messages"])
        print('task: ', state['messages'][0].content)
        
        draft = enhanced_search(state['messages'][0].content)  
        print('draft: ', draft)
        
        return {
            "messages": [AIMessage(content=draft)]
        }
    
    class Reflection(BaseModel):
        missing: str = Field(description="Critique of what is missing.")
        advisable: str = Field(description="Critique of what is helpful for better answer")
        superfluous: str = Field(description="Critique of what is superfluous")

    class Research(BaseModel):
        """Provide reflection and then follow up with search queries to improve the answer."""

        reflection: Reflection = Field(description="Your reflection on the initial answer.")
        search_queries: list[str] = Field(
            description="1-3 search queries for researching improvements to address the critique of your current answer."
        )
    
    def reflect(state: State):
        print("###### reflect ######")
        print('state: ', state["messages"])    
        print('draft: ', state["messages"][-1].content)
    
        reflection = []
        search_queries = []
        for attempt in range(5):
            chat = get_chat()
            structured_llm = chat.with_structured_output(Research, include_raw=True)
            
            info = structured_llm.invoke(state["messages"][-1].content)
            print(f'attempt: {attempt}, info: {info}')
                
            if not info['parsed'] == None:
                parsed_info = info['parsed']
                # print('reflection: ', parsed_info.reflection)                
                reflection = [parsed_info.reflection.missing, parsed_info.reflection.advisable]
                search_queries = parsed_info.search_queries
                
                print('reflection: ', parsed_info.reflection)            
                print('search_queries: ', search_queries)                
                break
        
        return {
            "messages": state["messages"],
            "reflection": reflection,
            "search_queries": search_queries
        }

    def revise_answer(state: State):   
        print("###### revise_answer ######")
        system = """Revise your previous answer using the new information. 
You should use the previous critique to add important information to your answer. provide the final answer with <result> tag. 
<critique>
{reflection}
</critique>

<information>
{content}
</information>"""
                    
        reflection_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
            
        content = []        
        if useEnhancedSearch:
            for q in state["search_queries"]:
                response = enhanced_search(q)     
                print(f'q: {q}, response: {response}')
                content.append(response)                   
        else:
            search = TavilySearchResults(k=2)
            for q in state["search_queries"]:
                response = search.invoke(q)     
                for r in response:
                    content.append(r['content'])     

        chat = get_chat()
        reflect = reflection_prompt | chat
            
        messages = state["messages"]
        cls_map = {"ai": HumanMessage, "human": AIMessage}
        translated = [messages[0]] + [
            cls_map[msg.type](content=msg.content) for msg in messages[1:]
        ]
        print('translated: ', translated)     
           
        res = reflect.invoke(
            {
                "messages": translated,
                "reflection": state["reflection"],
                "content": content
            }
        )    
                                
        response = HumanMessage(content=res.content[res.content.find('<result>')+8:len(res.content)-9])
        print('response: ', response)
                
        revision_number = state["revision_number"] if state.get("revision_number") is not None else 1
        return {
            "messages": [response], 
            "revision_number": revision_number + 1
        }
    
    MAX_REVISIONS = 1
    def should_continue(state: State, config):
        print("###### should_continue ######")
        max_revisions = config.get("configurable", {}).get("max_revisions", MAX_REVISIONS)
        print("max_revisions: ", max_revisions)
            
        if state["revision_number"] > max_revisions:
            return "end"
        return "contine"

    def buildKnowledgeGuru():    
        workflow = StateGraph(State)

        workflow.add_node("generate", generate)
        workflow.add_node("reflect", reflect)
        workflow.add_node("revise_answer", revise_answer)

        workflow.set_entry_point("generate")

        workflow.add_conditional_edges(
            "revise_answer", 
            should_continue, 
            {
                "end": END, 
                "contine": "reflect"}
        )

        workflow.add_edge("generate", "reflect")
        workflow.add_edge("reflect", "revise_answer")
        
        app = workflow.compile()
        
        return app
    
    app = buildKnowledgeGuru()
        
    isTyping(connectionId, requestId)    
    inputs = [HumanMessage(content=query)]
    config = {
        "recursion_limit": 50,
        "max_revisions": MAX_REVISIONS
    }
    
    for output in app.stream({"messages": inputs}, config):   
        for key, value in output.items():
            print(f"Finished: {key}")
            #print("value: ", value)
            
    print('value: ', value)
        
    readStreamMsg(connectionId, requestId, value["messages"][-1].content)
    
    return value["messages"][-1].content

#########################################################
def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)
    function_type = jsonBody['function_type']  # conversation type
    print('Function Type: ', function_type)
    
    print('initiate....')
    global reference_docs
    reference_docs = []
    
    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)

    global enableReference, code_type
    global map_chain, memory_chain, debugMessageMode, allowDualSearch
    
    if function_type == 'dual-search':
        allowDualSearch = 'true'
    elif function_type == 'code-generation-python':
        code_type = 'py'
    elif function_type == 'code-generation-nodejs':
        code_type = 'js'

    # Multi-LLM
    profile = LLM_for_chat[selected_chat]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_chat: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    # print('profile: ', profile)
    
    chat = get_chat()    
    bedrock_embedding = get_embedding()

    # allocate memory
    if userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        allowTime = getAllowTime()
        load_chat_history(userId, allowTime)
 
    start = int(time.time())    

    msg = ""
    reference = ""
    speech_uri = ""
    token_counter_input = 0
    time_for_inference = 0
    history_length = 0
    isControlMsg = False
    
    if type == 'text' and body[:11] == 'list models':
        isControlMsg = True
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:           
        text = body
        print('query: ', text)  
        querySize = len(text)
        textCount = len(text.split())
        print(f"query size: {querySize}, words: {textCount}")
        
        if type == 'text':
            if text == 'enableReference':
                enableReference = 'true'
                isControlMsg = True
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                isControlMsg = True
                msg  = "Reference is disabled"
            elif text == 'useOpenSearch':
                rag_type = 'opensearch'
                isControlMsg = True
                msg  = "OpenSearch is selected for Knowledge Database"
            elif text == 'useKendra':
                isControlMsg = True
                rag_type = 'kendra'
                msg  = "Kendra is selected for Knowledge Database"
            elif text == 'enableDebug':
                isControlMsg = True
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                isControlMsg = True
                debugMessageMode = 'false'
                msg  = "Debug messages will not be delivered to the client."
            elif text == 'enableDualSearch':
                isControlMsg = True
                allowDualSearch = 'true'
                msg  = "Dual Search is enabled"
            elif text == 'disableDualSearch':
                isControlMsg = True
                allowDualSearch = 'false'
                msg  = "Dual Search is disabled"

            elif text == 'clearMemory':
                memory_chain.clear()
                map_chain[userId] = memory_chain
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:                       
                if conv_type == 'normal' or conv_type == 'funny':      # normal
                    msg = general_conversation(connectionId, requestId, chat, text)        
                              
                elif conv_type == 'agent-executor':                    
                    msg = run_agent_executor(connectionId, requestId, text)
                    if reference_docs:
                        reference = get_references_for_agent(reference_docs)      
                        
                elif conv_type == 'agent-executor-chat':
                    #revised_question = revise_question(connectionId, requestId, chat, text)     
                    #print('revised_question: ', revised_question)  
                    msg = run_agent_executor(connectionId, requestId, text)  
                    if reference_docs:
                        reference = get_references_for_agent(reference_docs)      
                                                      
                elif conv_type == 'agent-reflection':  # reflection
                    msg = run_reflection_agent(connectionId, requestId, text)     
                    
                elif conv_type == 'agent-knowledge-guru':  # knowledge guru
                    msg = run_knowledge_guru(connectionId, requestId, text)      
                
                elif conv_type == 'qa':   # RAG
                    print(f'rag_type: {rag_type}')
                    msg, reference = get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding, rag_type)
                    
                elif conv_type == "qa-kb":
                    msg, reference = get_answer_using_knowledge_base(chat, text, connectionId, requestId)                
                elif conv_type == "qa-kb-chat":
                    revised_question = revise_question(connectionId, requestId, chat, text)     
                    print('revised_question: ', revised_question)      
                    msg, reference = get_answer_using_knowledge_base(chat, revised_question, connectionId, requestId)                
                
                elif conv_type == "prompt-flow":
                    msg = run_prompt_flow(text, connectionId, requestId)
                elif conv_type == "prompt-flow-chat":
                    revised_question = revise_question(connectionId, requestId, chat, text)     
                    print('revised_question: ', revised_question)                    
                    msg = run_prompt_flow(revised_question, connectionId, requestId)
                
                elif conv_type == "rag-prompt-flow":
                    msg = run_RAG_prompt_flow(text, connectionId, requestId)
                
                elif conv_type == "bedrock-agent":
                    msg = run_bedrock_agent(text, connectionId, requestId)
                
                elif conv_type == "translation":
                    msg = translate_text(chat, text) 
                elif conv_type == "grammar":
                    msg = check_grammer(chat, text)  
                elif conv_type == "sentiment":
                    msg = extract_sentiment(chat, text)
                elif conv_type == "extraction": # infomation extraction
                    msg = extract_information(chat, text)  
                elif conv_type == "pii":
                    msg = remove_pii(chat, text)   
                elif conv_type == "step-by-step":
                    msg = do_step_by_step(chat, text)  
                elif conv_type == "timestamp-extraction":
                    msg = extract_timestamp(chat, text)  
                else:
                    msg = general_conversation(connectionId, requestId, chat, text) 
                    
                # token counter
                if debugMessageMode=='true':
                    token_counter_input = chat.get_num_tokens(text)
                    token_counter_output = chat.get_num_tokens(msg)
                    print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                
                memory_chain.chat_memory.add_user_message(text)  # append new diaglog
                memory_chain.chat_memory.add_ai_message(msg)
                
        elif type == 'code':
            msg, reference = get_code_using_RAG(chat, text, code_type, connectionId, requestId, bedrock_embedding)  
            
            # token counter
            if debugMessageMode=='true':
                token_counter_input = chat.get_num_tokens(text)
                token_counter_output = chat.get_num_tokens(msg)
                print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                
        elif type == 'document':
            isTyping(connectionId, requestId)
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)
            
            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
            
            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
                
            elif file_type == 'py' or file_type == 'js':
                s3r = boto3.resource("s3")
                doc = s3r.Object(s3_bucket, s3_prefix+'/'+object)
                
                contents = doc.get()['Body'].read().decode('utf-8')
                
                #contents = load_code(file_type, object)                
                                
                msg = summary_of_code(chat, contents, file_type)                  
                
            elif file_type == 'png' or file_type == 'jpeg' or file_type == 'jpg':
                print('multimodal: ', object)
                
                s3_client = boto3.client('s3') 
                    
                image_obj = s3_client.get_object(Bucket=s3_bucket, Key=s3_prefix+'/'+object)
                # print('image_obj: ', image_obj)
                
                image_content = image_obj['Body'].read()
                img = Image.open(BytesIO(image_content))
                
                width, height = img.size 
                print(f"width: {width}, height: {height}, size: {width*height}")
                
                isResized = False
                while(width*height > 5242880):                    
                    width = int(width/2)
                    height = int(height/2)
                    isResized = True
                    print(f"width: {width}, height: {height}, size: {width*height}")
                
                if isResized:
                    img = img.resize((width, height))
                
                buffer = BytesIO()
                img.save(buffer, format="PNG")
                img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                
                command  = jsonBody['command']
                print('command: ', command)
                
                # verify the image
                msg = use_multimodal(img_base64, command)       
                
                # extract text from the image
                text = extract_text(chat, img_base64)
                extracted_text = text[text.find('<result>')+8:len(text)-9] # remove <result> tag
                print('extracted_text: ', extracted_text)
                if len(extracted_text)>10:
                    msg = msg + f"\n\n[추출된 Text]\n{extracted_text}\n"
                
                memory_chain.chat_memory.add_user_message(f"{object}에서 텍스트를 추출하세요.")
                memory_chain.chat_memory.add_ai_message(extracted_text)
            
            else:
                msg = "uploaded file: "+object
                                                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)

        # translation and tts(text to speech)
        if type=='text' and isControlMsg==False:        
            if (conv_type=='qa' or  conv_type == "normal"):
                if isKorean(msg)==False :
                    translated_msg = traslation_to_korean(chat=chat, text=msg)
                    print('translated_msg: ', translated_msg)
                    msg = msg+'\n[한국어]\n'+translated_msg
                    
                    if speech_generation == 'true': # generate mp3 file
                        speech_uri = get_text_speech(path=path, speech_prefix=speech_prefix, bucket=s3_bucket, msg=translated_msg)
                        print('speech_uri: ', speech_uri)                      
                else:
                    if speech_generation == 'true': # generate mp3 file
                        speech_uri = get_text_speech(path=path, speech_prefix=speech_prefix, bucket=s3_bucket, msg=msg)
                        print('speech_uri: ', speech_uri)  
           
        if type == 'code':  # code summary
            if reference: # Summarize the generated code 
                generated_code = msg[msg.find('<result>')+9:len(msg)-10]
                generated_code_summary = summary_of_code(chat, generated_code, code_type)    
                msg += f'\n\n[생성된 코드 설명]\n{generated_code_summary}'
                msg = msg.replace('\n\n\n', '\n\n') 
                    
                sendResultMessage(connectionId, requestId, msg+reference)

        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if speech_uri:
        speech = '\n' + f'<a href={speech_uri} target=_blank>{"[결과 읽어주기 (mp3)]"}</a>'                
        sendResultMessage(connectionId, requestId, msg+reference+speech)

        if conv_type=='qa' and debugMessageMode=='true' and reference:
            statusMsg = f"\n[통계]\nRegion: {bedrock_region}\nQuestion: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            statusMsg = statusMsg + f"RAG: {str(relevant_length)}자 / {token_counter_relevant_docs}토큰 ({number_of_relevant_docs})\n"

            statusMsg = statusMsg + f"Time(초): "            
            if time_for_revise != 0:
                statusMsg = statusMsg + f"{time_for_revise:.2f}(Revise), "
            if time_for_rag != 0:
                statusMsg = statusMsg + f"{time_for_rag:.2f}(RAG), "
            if time_for_priority_search != 0:
                statusMsg = statusMsg + f"{time_for_priority_search:.2f}(Priority) "
            if time_for_inference != 0:
                statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
            statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
            if time_for_rag_inference != 0:
                statusMsg = statusMsg + f"\nRAG-Detail: {time_for_rag_inference:.2f}(Inference(KOR)), "
            if time_for_rag_question_translation != 0:
                statusMsg = statusMsg + f"{time_for_rag_question_translation:.2f}(Question(ENG)), "
            if time_for_rag_2nd_inference != 0:
                statusMsg = statusMsg + f"{time_for_rag_2nd_inference:.2f}(Inference(ENG)), "
            if time_for_rag_translation != 0:
                statusMsg = statusMsg + f"{time_for_rag_translation:.2f}(Doc Translation), "

            sendResultMessage(connectionId, requestId, msg+reference+speech+statusMsg)
        elif debugMessageMode=='true': # other cases
            statusMsg = f"\n[통계]\nRegion: {bedrock_region}\n"
            if token_counter_input:
                statusMsg = statusMsg + f"Question: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            
            if history_length:
                statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            
            statusMsg = statusMsg + f"Time(초): "            
            if time_for_inference != 0:
                statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
            statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
            sendResultMessage(connectionId, requestId, msg+speech+statusMsg)

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
